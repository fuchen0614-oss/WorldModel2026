#!/usr/bin/env python3
"""Render a diagnostic PNG from simple native PowerPoint shapes.

This is intentionally approximate. It supports the editable primitives used by
refine_user_ppts.py and is used only for iterative visual QA when PowerPoint or
LibreOffice is unavailable on the workspace node.
"""

from __future__ import annotations

import argparse
import math
import sys
from pathlib import Path


ROOT = Path(__file__).resolve().parents[2]
sys.path.insert(0, str(ROOT / ".tools" / "python-pptx"))

from PIL import Image, ImageDraw, ImageFont  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402


FONT_REGULAR = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf")
FONT_BOLD = Path("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf")
FONT_REGULAR_FALLBACK = Path("/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf")
FONT_BOLD_FALLBACK = Path("/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf")


def color_to_tuple(color, default=(255, 255, 255)):
    try:
        value = color.rgb
        if value is not None:
            return tuple(value)
    except Exception:
        pass
    return default


def get_font(size: int, bold: bool):
    candidates = (
        (FONT_BOLD, FONT_BOLD_FALLBACK)
        if bold
        else (FONT_REGULAR, FONT_REGULAR_FALLBACK)
    )
    for candidate in candidates:
        if candidate.exists():
            return ImageFont.truetype(str(candidate), max(7, size))
    return ImageFont.load_default()


def wrap_text(draw, text, font, max_width):
    lines = []
    for explicit_line in (text or "").splitlines() or [""]:
        if not explicit_line:
            lines.append("")
            continue
        words = explicit_line.split(" ")
        current = ""
        for word in words:
            candidate = word if not current else f"{current} {word}"
            if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
                current = candidate
            else:
                if current:
                    lines.append(current)
                current = word
        lines.append(current)
    return lines


def draw_text(draw, shape, scale, x, y, w, h):
    if not getattr(shape, "has_text_frame", False):
        return
    text = shape.text.strip()
    if not text:
        return
    p = shape.text_frame.paragraphs[0]
    run = p.runs[0] if p.runs else None
    pt = 12
    bold = False
    color = (31, 41, 55)
    if run is not None:
        try:
            if run.font.size:
                pt = run.font.size.pt
        except Exception:
            pass
        bold = bool(run.font.bold)
        color = color_to_tuple(run.font.color, color)
    font = get_font(round(pt * 12700 * scale), bold)
    margin = max(2, round(0.04 * scale))
    lines = wrap_text(draw, text, font, max(4, w - 2 * margin))
    line_h = draw.textbbox((0, 0), "Ag", font=font)[3] + max(1, round(0.01 * scale))
    block_h = line_h * len(lines)
    ty = y + max(margin, (h - block_h) / 2)
    align = str(p.alignment or "")
    for line in lines:
        box = draw.textbbox((0, 0), line, font=font)
        tw = box[2] - box[0]
        if "LEFT" in align:
            tx = x + margin
        elif "RIGHT" in align:
            tx = x + w - margin - tw
        else:
            tx = x + (w - tw) / 2
        draw.text((tx, ty), line, fill=color, font=font)
        ty += line_h


def add_arrow(draw, p1, p2, color, width):
    draw.line([p1, p2], fill=color, width=width)
    angle = math.atan2(p2[1] - p1[1], p2[0] - p1[0])
    size = max(5, width * 3)
    pts = [
        p2,
        (
            p2[0] - size * math.cos(angle - math.pi / 6),
            p2[1] - size * math.sin(angle - math.pi / 6),
        ),
        (
            p2[0] - size * math.cos(angle + math.pi / 6),
            p2[1] - size * math.sin(angle + math.pi / 6),
        ),
    ]
    draw.polygon(pts, fill=color)


def render(pptx_path: Path, png_path: Path, slide_index: int = 0, width: int = 2400):
    prs = Presentation(pptx_path)
    slide = prs.slides[slide_index]
    scale = width / prs.slide_width
    height = round(prs.slide_height * scale)
    canvas = Image.new("RGB", (width, height), "white")
    draw = ImageDraw.Draw(canvas)

    def px(v):
        return round(v * scale)

    for shape in slide.shapes:
        x, y, w, h = map(px, (shape.left, shape.top, shape.width, shape.height))
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            color = color_to_tuple(shape.line.color, (70, 80, 95))
            lw = max(1, round((shape.line.width.pt if shape.line.width else 1.2) * 12700 * scale))
            add_arrow(draw, (x, y), (x + w, y + h), color, lw)
            continue

        if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            draw_text(draw, shape, scale, x, y, w, h)
            continue

        try:
            fill = color_to_tuple(shape.fill.fore_color, (255, 255, 255))
        except Exception:
            fill = (255, 255, 255)
        line = color_to_tuple(shape.line.color, (180, 185, 195))
        lw = max(1, round((shape.line.width.pt if shape.line.width else 1.0) * 12700 * scale))
        kind = str(getattr(shape, "auto_shape_type", ""))
        box = [x, y, x + w, y + h]
        if "OVAL" in kind:
            draw.ellipse(box, fill=fill, outline=line, width=lw)
        elif "ISOSCELES_TRIANGLE" in kind:
            draw.polygon([(x + w / 2, y), (x + w, y + h), (x, y + h)], fill=fill, outline=line)
        elif "ROUNDED_RECTANGLE" in kind:
            draw.rounded_rectangle(box, radius=max(4, min(w, h) // 10), fill=fill, outline=line, width=lw)
        else:
            draw.rectangle(box, fill=fill, outline=line, width=lw)
        draw_text(draw, shape, scale, x, y, w, h)

    png_path.parent.mkdir(parents=True, exist_ok=True)
    canvas.save(png_path)


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("pptx", type=Path)
    parser.add_argument("png", type=Path)
    parser.add_argument("--slide", type=int, default=1)
    parser.add_argument("--width", type=int, default=2400)
    args = parser.parse_args()
    render(args.pptx, args.png, args.slide - 1, args.width)
    print(args.png)


if __name__ == "__main__":
    main()
