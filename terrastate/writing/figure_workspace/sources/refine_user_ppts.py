#!/usr/bin/env python3
"""Rebuild only slide 1 of the user's two editable figure decks.

The deck names reflect the user's working files rather than final paper numbering:

* 示例/fig1.pptx currently contains the four-region TerraState method figure.
* 示例/fig2.pptx currently contains the three-region conceptual/evidence figure.

All later slides are preserved byte-for-byte by transplanting only slide1.xml
back into the original OOXML package. No external paper artwork is copied.
Missing project imagery is represented by native, replaceable PowerPoint slots.
"""

from __future__ import annotations

import shutil
import sys
import tempfile
import zipfile
from pathlib import Path


ROOT = Path(__file__).resolve().parents[2]
TOOL = ROOT / ".tools" / "python-pptx"
sys.path.insert(0, str(TOOL))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.dml import MSO_LINE_DASH_STYLE  # noqa: E402
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN  # noqa: E402
from pptx.oxml.xmlchemy import OxmlElement  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


EXAMPLE = ROOT / "示例"
OUTDIR = EXAMPLE / "codex_drafts"
OUTDIR.mkdir(parents=True, exist_ok=True)

FONT = "Arial"

INK = "1F2937"
MUTED = "667085"
LINE = "98A2B3"
WHITE = "FFFFFF"
PANEL = "FBFCFE"
BLUE = "2F6F9F"
BLUE_LIGHT = "EAF3F8"
TEAL = "238B7B"
TEAL_LIGHT = "EAF6F3"
PURPLE = "7256A8"
PURPLE_LIGHT = "F2EFF9"
ORANGE = "C96A16"
ORANGE_LIGHT = "FFF3E8"
GREEN = "467A55"
GREEN_LIGHT = "EDF5EE"
RED = "B54745"
GRAY_LIGHT = "F2F4F7"
PLACEHOLDER = "F7F8FA"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_name(shape, name: str):
    shape.name = name
    return shape


def remove_all_shapes(slide) -> None:
    tree = slide.shapes._spTree
    for shape in list(slide.shapes):
        tree.remove(shape._element)


def set_slide_white(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(WHITE)


def set_cell_text(
    shape,
    text: str,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.03,
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(
    slide,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    size: float = 18,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
    valign=MSO_ANCHOR.MIDDLE,
    margin: float = 0.02,
):
    shape = set_name(
        slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)), name
    )
    set_cell_text(shape, text, size, color, bold, align, valign, margin)
    return shape


def add_box(
    slide,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = WHITE,
    line: str = LINE,
    width: float = 1.2,
    radius=True,
    dash: MSO_LINE_DASH_STYLE | None = None,
    text: str | None = None,
    text_size: float = 18,
    text_color: str = INK,
    bold: bool = False,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = set_name(
        slide.shapes.add_shape(
            kind, Inches(x), Inches(y), Inches(w), Inches(h)
        ),
        name,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    shape.line.width = Pt(width)
    if dash is not None:
        shape.line.dash_style = dash
    if text is not None:
        set_cell_text(shape, text, text_size, text_color, bold)
    return shape


def add_panel(
    slide,
    prefix: str,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    tint: str,
    accent: str,
):
    add_box(
        slide,
        f"{prefix}_panel",
        x,
        y,
        w,
        h,
        PANEL,
        "D0D5DD",
        1.0,
        True,
    )
    add_box(
        slide,
        f"{prefix}_header",
        x,
        y,
        w,
        0.46,
        tint,
        accent,
        1.1,
        True,
        text=title,
        text_size=18,
        text_color=INK,
        bold=True,
    )


def add_arrowhead(connector, color: str = INK, size: str = "med") -> None:
    ln = connector._element.spPr.ln
    for child in list(ln):
        if child.tag.endswith("tailEnd"):
            ln.remove(child)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", size)
    tail.set("len", size)
    ln.append(tail)


def add_line(
    slide,
    name: str,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = INK,
    width: float = 1.4,
    arrow: bool = True,
    dash: MSO_LINE_DASH_STYLE | None = None,
):
    line = set_name(
        slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        ),
        name,
    )
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dash is not None:
        line.line.dash_style = dash
    if arrow:
        add_arrowhead(line, color)
    return line


def add_elbow(
    slide,
    prefix: str,
    points: list[tuple[float, float]],
    color: str = INK,
    width: float = 1.3,
    dash: MSO_LINE_DASH_STYLE | None = None,
):
    for index, ((x1, y1), (x2, y2)) in enumerate(zip(points, points[1:]), start=1):
        add_line(
            slide,
            f"{prefix}_{index}",
            x1,
            y1,
            x2,
            y2,
            color,
            width,
            arrow=index == len(points) - 1,
            dash=dash,
        )


def add_image_slot(
    slide,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    hint: str,
    accent: str = BLUE,
    fill: str = PLACEHOLDER,
    frames: int = 1,
):
    add_box(
        slide,
        name,
        x,
        y,
        w,
        h,
        fill,
        accent,
        1.0,
        True,
        MSO_LINE_DASH_STYLE.DASH,
    )
    top = 0.12
    usable_h = max(0.25, h - 0.54)
    if frames == 1:
        add_box(
            slide,
            f"{name}_visual",
            x + 0.12,
            y + top,
            w - 0.24,
            usable_h,
            WHITE,
            "D0D5DD",
            0.8,
            False,
        )
    else:
        card_w = (w - 0.30 - 0.08 * (frames - 1)) / frames
        for idx in range(frames):
            add_box(
                slide,
                f"{name}_frame_{idx+1}",
                x + 0.15 + idx * (card_w + 0.08),
                y + top + idx * 0.02,
                card_w,
                usable_h - idx * 0.02,
                WHITE,
                "D0D5DD",
                0.8,
                False,
            )
    add_text(
        slide,
        f"{name}_title",
        x + 0.08,
        y + h - 0.37,
        w - 0.16,
        0.18,
        title,
        12.5,
        INK,
        True,
    )
    add_text(
        slide,
        f"{name}_hint",
        x + 0.08,
        y + h - 0.20,
        w - 0.16,
        0.15,
        hint,
        11.0,
        MUTED,
        False,
    )


def add_token_wall(
    slide,
    prefix: str,
    x: float,
    y: float,
    cols: int = 4,
    rows: int = 4,
    cell: float = 0.12,
    gap: float = 0.035,
    colors: tuple[str, ...] = (BLUE, TEAL, PURPLE),
):
    for c in range(cols):
        for r in range(rows):
            add_box(
                slide,
                f"{prefix}_{c}_{r}",
                x + c * (cell + gap),
                y + r * (cell + gap),
                cell,
                cell,
                colors[c % len(colors)],
                WHITE,
                0.3,
                False,
            )


def add_state_grid(
    slide,
    prefix: str,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    accent: str = TEAL,
    changed: bool = False,
):
    add_box(slide, prefix, x, y, w, h, TEAL_LIGHT, accent, 1.3, True)
    grid_x, grid_y = x + 0.18, y + 0.14
    cell = min((w - 0.36) / 4.0, (h - 0.55) / 4.0)
    fills = [TEAL, "57A89B", "8BC2B9", "C6E2DD"]
    for c in range(4):
        for r in range(4):
            idx = (c + r + (2 if changed else 0)) % len(fills)
            add_box(
                slide,
                f"{prefix}_cell_{c}_{r}",
                grid_x + c * cell,
                grid_y + r * cell,
                cell - 0.015,
                cell - 0.015,
                fills[idx],
                WHITE,
                0.25,
                False,
            )
    add_text(
        slide,
        f"{prefix}_label",
        x + 0.03,
        y + h - 0.32,
        w - 0.06,
        0.25,
        label,
        14,
        INK,
        True,
    )


def add_weather_strip(
    slide,
    prefix: str,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    compact: bool = False,
):
    add_box(slide, prefix, x, y, w, h, ORANGE_LIGHT, ORANGE, 1.1, True)
    rows = [("T", ORANGE), ("P", BLUE), ("R", PURPLE)]
    row_h = (h - 0.35) / 3
    for idx, (tag, color) in enumerate(rows):
        yy = y + 0.07 + idx * row_h
        add_text(
            slide,
            f"{prefix}_{tag}_label",
            x + 0.06,
            yy,
            0.18,
            row_h,
            tag,
            10.5,
            color,
            True,
        )
        points = [
            (x + 0.28, yy + row_h * 0.62),
            (x + w * 0.38, yy + row_h * (0.33 if idx != 1 else 0.72)),
            (x + w * 0.58, yy + row_h * (0.72 if idx == 0 else 0.36)),
            (x + w - 0.10, yy + row_h * (0.42 if idx != 2 else 0.66)),
        ]
        for j, ((x1, y1), (x2, y2)) in enumerate(zip(points, points[1:]), start=1):
            add_line(
                slide,
                f"{prefix}_{tag}_segment_{j}",
                x1,
                y1,
                x2,
                y2,
                color,
                0.8,
                False,
            )
    add_text(
        slide,
        f"{prefix}_label",
        x + 0.03,
        y + h - 0.25,
        w - 0.06,
        0.19,
        label,
        11.5 if compact else 12.0,
        INK,
        True,
    )


def add_small_map_pair(
    slide,
    prefix: str,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str = "geographic context g",
):
    add_box(slide, prefix, x, y, w, h, GREEN_LIGHT, GREEN, 1.0, True)
    half = (w - 0.20) / 2
    add_box(
        slide,
        f"{prefix}_landcover",
        x + 0.07,
        y + 0.08,
        half,
        h - 0.34,
        "D8E8C9",
        "A8B79A",
        0.7,
        False,
    )
    # Land-cover patches.
    for idx, (dx, dy, fw, fh, color) in enumerate(
        [
            (0.05, 0.05, 0.20, 0.14, GREEN),
            (0.27, 0.04, 0.18, 0.22, "D5A848"),
            (0.08, 0.23, 0.30, 0.14, "7AA66C"),
        ],
        start=1,
    ):
        add_box(
            slide,
            f"{prefix}_patch_{idx}",
            x + 0.07 + dx,
            y + 0.08 + dy,
            min(fw, half - dx - 0.02),
            min(fh, h - 0.36 - dy),
            color,
            color,
            0,
            False,
        )
    add_box(
        slide,
        f"{prefix}_dem",
        x + 0.13 + half,
        y + 0.08,
        half,
        h - 0.34,
        "E7E9EC",
        "A8ADB5",
        0.7,
        False,
    )
    # DEM contour lines.
    for idx in range(3):
        add_line(
            slide,
            f"{prefix}_contour_{idx}",
            x + 0.17 + half,
            y + 0.15 + idx * 0.11,
            x + 0.09 + 2 * half,
            y + 0.20 + idx * 0.08,
            MUTED,
            0.6,
            False,
        )
    add_text(
        slide,
        f"{prefix}_label",
        x + 0.03,
        y + h - 0.24,
        w - 0.06,
        0.18,
        label,
        11.5,
        INK,
        True,
    )


def add_clock_chip(
    slide,
    prefix: str,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str = "h",
):
    add_box(
        slide,
        prefix,
        x,
        y,
        w,
        h,
        GRAY_LIGHT,
        MUTED,
        1.0,
        True,
        text=f"horizon\n{value}",
        text_size=11,
        bold=True,
    )


def add_plus(slide, prefix: str, x: float, y: float, diameter: float = 0.42):
    shape = set_name(
        slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x),
            Inches(y),
            Inches(diameter),
            Inches(diameter),
        ),
        prefix,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(WHITE)
    shape.line.color.rgb = rgb(ORANGE)
    shape.line.width = Pt(1.8)
    set_cell_text(shape, "+", 18, ORANGE, True)
    return shape


def add_intervention_port(
    slide,
    prefix: str,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    subtitle: str,
    color: str,
):
    add_box(
        slide,
        prefix,
        x,
        y,
        w,
        h,
        WHITE,
        color,
        1.2,
        True,
        MSO_LINE_DASH_STYLE.DASH,
    )
    add_text(
        slide,
        f"{prefix}_title",
        x + 0.05,
        y + 0.02,
        w - 0.10,
        h * 0.50,
        title,
        12.0,
        color,
        True,
    )
    add_text(
        slide,
        f"{prefix}_subtitle",
        x + 0.05,
        y + h * 0.46,
        w - 0.10,
        h * 0.46,
        subtitle,
        11.0,
        MUTED,
        False,
    )


def build_method_slide(slide) -> None:
    """Four-region, inference-only TerraState method architecture."""

    remove_all_shapes(slide)
    set_slide_white(slide)

    y, h = 0.16, 7.16
    gap = 0.08
    x_a, w_a = 0.16, 2.34
    x_b, w_b = x_a + w_a + gap, 3.10
    x_c, w_c = x_b + w_b + gap, 4.16
    x_d, w_d = x_c + w_c + gap, 3.19

    add_panel(
        slide,
        "A",
        x_a,
        y,
        w_a,
        h,
        "(a) Multimodal context",
        BLUE_LIGHT,
        BLUE,
    )
    add_panel(
        slide,
        "B",
        x_b,
        y,
        w_b,
        h,
        "(b) History encoding & state construction",
        TEAL_LIGHT,
        TEAL,
    )
    add_panel(
        slide,
        "C",
        x_c,
        y,
        w_c,
        h,
        "(c) Weather-conditioned shared dynamics",
        PURPLE_LIGHT,
        PURPLE,
    )
    add_panel(
        slide,
        "D",
        x_d,
        y,
        w_d,
        h,
        "(d) State readout & forecast",
        ORANGE_LIGHT,
        ORANGE,
    )

    # A: project-image slots.
    add_image_slot(
        slide,
        "A1_EO_history_slot",
        x_a + 0.16,
        0.82,
        w_a - 0.32,
        1.22,
        "Historical Earth observations",
        "replace with 3–4 aligned EO / NDVI frames",
        BLUE,
        frames=3,
    )
    add_image_slot(
        slide,
        "A2_environment_slot",
        x_a + 0.16,
        2.20,
        w_a - 0.32,
        1.04,
        "Historical environmental context",
        "cloud mask + past meteorology",
        BLUE,
        frames=2,
    )
    add_small_map_pair(
        slide,
        "A3_geography_slot",
        x_a + 0.16,
        3.40,
        w_a - 0.32,
        1.02,
        "static geography: land cover + DEM",
    )
    add_weather_strip(
        slide,
        "A4_future_weather",
        x_a + 0.16,
        4.62,
        w_a - 0.32,
        1.55,
        "future meteorological forcing u(t:t+h)",
    )
    add_text(
        slide,
        "A_image_note",
        x_a + 0.19,
        6.34,
        w_a - 0.38,
        0.55,
        "Image slots preserve space for provenance-frozen project data.",
        10.5,
        MUTED,
        False,
        PP_ALIGN.LEFT,
    )

    # B: q branches to context forecast and predictive state.
    add_box(
        slide,
        "B1_history_encoder",
        x_b + 0.30,
        1.05,
        2.50,
        1.12,
        BLUE_LIGHT,
        BLUE,
        1.5,
        True,
    )
    add_token_wall(slide, "B1_input_tokens", x_b + 0.47, 1.36, 3, 3, 0.10, 0.03)
    add_text(
        slide,
        "B1_encoder_label",
        x_b + 1.07,
        1.20,
        1.52,
        0.50,
        "History encoder q",
        17,
        INK,
        True,
    )
    add_text(
        slide,
        "B1_encoder_subtitle",
        x_b + 1.07,
        1.66,
        1.52,
        0.28,
        "cloud-aware context encoding",
        10.5,
        MUTED,
    )

    add_text(
        slide,
        "B2_branch_title",
        x_b + 0.22,
        2.42,
        2.66,
        0.28,
        "State-construction branch",
        13,
        TEAL,
        True,
    )
    add_token_wall(slide, "B2_context_features", x_b + 0.28, 2.90, 3, 4, 0.12, 0.035)
    add_text(
        slide,
        "B2_context_features_label",
        x_b + 0.16,
        3.56,
        0.92,
        0.30,
        "context\nfeatures",
        10.5,
        MUTED,
        False,
    )
    add_box(
        slide,
        "B2_projector",
        x_b + 1.25,
        3.04,
        0.68,
        0.66,
        TEAL_LIGHT,
        TEAL,
        1.2,
        True,
        text="P\nprojector",
        text_size=12,
        bold=True,
    )
    add_state_grid(
        slide,
        "B2_state_zt",
        x_b + 2.12,
        2.77,
        0.78,
        1.24,
        "z_t",
        TEAL,
    )
    add_line(
        slide,
        "B2_features_to_P",
        x_b + 0.90,
        3.27,
        x_b + 1.25,
        3.27,
        TEAL,
        1.4,
    )
    add_line(
        slide,
        "B2_P_to_state",
        x_b + 1.93,
        3.27,
        x_b + 2.12,
        3.27,
        TEAL,
        1.4,
    )
    add_elbow(
        slide,
        "B1_to_state_branch",
        [
            (x_b + 1.15, 2.17),
            (x_b + 1.15, 2.58),
            (x_b + 0.66, 2.58),
            (x_b + 0.66, 2.88),
        ],
        TEAL,
        1.3,
    )

    add_box(
        slide,
        "B3_context_forecast",
        x_b + 0.30,
        4.70,
        2.50,
        1.08,
        BLUE_LIGHT,
        BLUE,
        1.2,
        True,
    )
    for idx, horizon in enumerate(("h=5", "h=10", "h=20")):
        add_box(
            slide,
            f"B3_forecast_slot_{idx+1}",
            x_b + 0.48 + idx * 0.67,
            4.92,
            0.54,
            0.45,
            WHITE,
            "B9C9D7",
            0.7,
            False,
            text=horizon,
            text_size=8.5,
            text_color=MUTED,
        )
    add_text(
        slide,
        "B3_context_forecast_label",
        x_b + 0.34,
        5.43,
        2.42,
        0.25,
        "Context-only forecasts b(1:H)",
        12.5,
        BLUE,
        True,
    )
    add_elbow(
        slide,
        "B1_to_context_branch",
        [
            (x_b + 1.92, 2.17),
            (x_b + 1.92, 4.46),
            (x_b + 1.55, 4.46),
            (x_b + 1.55, 4.70),
        ],
        BLUE,
        1.2,
    )
    add_text(
        slide,
        "B_history_only_note",
        x_b + 0.34,
        6.28,
        2.42,
        0.50,
        "Both branches read historical information only.",
        10.5,
        MUTED,
        False,
    )

    # Cross-panel historical-input arrow.
    add_line(
        slide,
        "A_history_to_B_q",
        x_a + w_a - 0.03,
        1.60,
        x_b + 0.30,
        1.60,
        BLUE,
        1.8,
    )

    # C1: conditions around the shared transition.
    add_weather_strip(
        slide,
        "C1_weather_condition",
        x_c + 0.22,
        0.88,
        1.95,
        1.12,
        "future weather u(t:t+h)",
        True,
    )
    add_small_map_pair(
        slide,
        "C1_geo_condition",
        x_c + 2.30,
        0.88,
        1.02,
        1.12,
        "geography g",
    )
    add_clock_chip(
        slide,
        "C1_horizon_condition",
        x_c + 3.44,
        0.88,
        0.50,
        1.12,
        "h",
    )

    add_intervention_port(
        slide,
        "C4_weather_intervention",
        x_c + 0.28,
        2.17,
        1.86,
        0.62,
        "Weather intervention",
        "actual / matched donor / normalized mean",
        PURPLE,
    )

    add_box(
        slide,
        "C2_shared_transition",
        x_c + 0.72,
        3.02,
        2.16,
        2.14,
        PURPLE_LIGHT,
        PURPLE,
        1.8,
        True,
    )
    add_text(
        slide,
        "C2_transition_title",
        x_c + 0.86,
        3.18,
        1.88,
        0.35,
        "Shared transition T",
        18,
        PURPLE,
        True,
    )
    add_text(
        slide,
        "C2_transition_subtitle",
        x_c + 0.88,
        4.56,
        1.84,
        0.34,
        "state–forcing interaction",
        11.5,
        MUTED,
        False,
    )
    add_token_wall(
        slide,
        "C2_state_tokens",
        x_c + 0.96,
        3.78,
        4,
        2,
        0.105,
        0.035,
        (TEAL, "57A89B"),
    )
    add_text(
        slide,
        "C2_interaction_mark",
        x_c + 1.55,
        3.78,
        0.46,
        0.36,
        "×",
        20,
        PURPLE,
        True,
    )
    add_token_wall(
        slide,
        "C2_weather_tokens",
        x_c + 2.02,
        3.78,
        4,
        2,
        0.105,
        0.035,
        (ORANGE, "E9A15D"),
    )

    add_state_grid(
        slide,
        "C3_state_future",
        x_c + 3.12,
        3.38,
        0.84,
        1.50,
        "z_(t+h)",
        TEAL,
        True,
    )
    add_line(
        slide,
        "B_state_to_C_transition",
        x_b + 2.90,
        3.37,
        x_c + 0.72,
        3.90,
        TEAL,
        1.8,
    )
    add_line(
        slide,
        "C_transition_to_future_state",
        x_c + 2.88,
        4.08,
        x_c + 3.12,
        4.08,
        TEAL,
        1.8,
    )
    add_line(
        slide,
        "C_weather_to_T",
        x_c + 1.20,
        2.00,
        x_c + 1.56,
        3.02,
        ORANGE,
        1.5,
    )
    add_line(
        slide,
        "C_geo_to_T",
        x_c + 2.82,
        2.00,
        x_c + 2.40,
        3.02,
        GREEN,
        1.3,
    )
    add_line(
        slide,
        "C_h_to_T",
        x_c + 3.69,
        2.00,
        x_c + 2.72,
        3.02,
        MUTED,
        1.3,
    )
    add_text(
        slide,
        "C_shared_note",
        x_c + 0.52,
        5.62,
        3.12,
        0.48,
        "One weather-conditioned transition is shared across queried horizons.",
        11,
        MUTED,
        False,
    )

    # D: readout, explicit state contribution, closure, and output slots.
    add_box(
        slide,
        "D1_state_readout",
        x_d + 0.16,
        2.22,
        0.70,
        1.54,
        GREEN_LIGHT,
        GREEN,
        1.5,
        True,
        text="O\nstate\nreadout",
        text_size=14,
        bold=True,
    )
    add_image_slot(
        slide,
        "D1_state_contribution_slot",
        x_d + 1.00,
        2.22,
        0.90,
        1.54,
        "state contribution r_h",
        "signed map slot",
        GREEN,
        frames=1,
    )
    add_plus(slide, "D2_forecast_fusion", x_d + 2.08, 2.74, 0.46)
    add_line(
        slide,
        "C_state_to_D_readout",
        x_c + 3.96,
        4.10,
        x_d + 0.16,
        2.99,
        TEAL,
        1.8,
    )
    add_line(
        slide,
        "D_readout_to_contribution",
        x_d + 0.86,
        2.99,
        x_d + 1.00,
        2.99,
        GREEN,
        1.6,
    )
    add_line(
        slide,
        "D_contribution_to_plus",
        x_d + 1.90,
        2.99,
        x_d + 2.08,
        2.99,
        ORANGE,
        1.6,
    )
    add_intervention_port(
        slide,
        "D4_state_path_intervention",
        x_d + 1.02,
        1.17,
        1.44,
        0.60,
        "State-path intervention",
        "remove state contribution",
        ORANGE,
    )
    add_line(
        slide,
        "D4_intervention_pointer",
        x_d + 1.72,
        1.77,
        x_d + 1.84,
        2.88,
        ORANGE,
        1.0,
        False,
        MSO_LINE_DASH_STYLE.DASH,
    )

    add_image_slot(
        slide,
        "D3_forecast_slot",
        x_d + 2.62,
        1.56,
        0.43,
        2.82,
        "ŷ_(t+h)",
        "forecast",
        ORANGE,
        frames=1,
    )
    add_line(
        slide,
        "D_plus_to_forecast",
        x_d + 2.54,
        2.97,
        x_d + 2.62,
        2.97,
        ORANGE,
        1.8,
    )
    add_text(
        slide,
        "D_closure_equation",
        x_d + 0.30,
        4.22,
        2.55,
        0.38,
        "ŷ_(t+h) = b_h + r_h",
        17,
        ORANGE,
        True,
    )
    for idx, horizon in enumerate(("h=5", "h=10", "h=20")):
        add_box(
            slide,
            f"D3_horizon_slot_{idx+1}",
            x_d + 0.20 + idx * 0.94,
            5.02,
            0.78,
            1.12,
            PLACEHOLDER,
            ORANGE,
            1.0,
            True,
            MSO_LINE_DASH_STYLE.DASH,
            text=f"{horizon}\nNDVI image slot",
            text_size=10.5,
            text_color=MUTED,
        )
    add_text(
        slide,
        "D3_output_group_label",
        x_d + 0.25,
        6.20,
        2.60,
        0.30,
        "Vegetation forecast",
        13,
        ORANGE,
        True,
    )

    # Context-only branch bypasses C and enters the addition.
    add_elbow(
        slide,
        "B_context_to_D_plus",
        [
            (x_b + 2.80, 5.22),
            (x_c + 3.92, 5.22),
            (x_d + 2.30, 5.22),
            (x_d + 2.30, 3.20),
        ],
        BLUE,
        1.25,
    )
    add_text(
        slide,
        "context_branch_label",
        x_c + 1.88,
        5.28,
        1.60,
        0.24,
        "context-only b_h",
        10.5,
        BLUE,
        True,
    )


def add_scene_placeholder(
    slide,
    prefix: str,
    x: float,
    y: float,
    w: float,
    h: float,
    label: str,
    color: str,
):
    add_box(
        slide,
        prefix,
        x,
        y,
        w,
        h,
        PLACEHOLDER,
        color,
        1.0,
        True,
        MSO_LINE_DASH_STYLE.DASH,
    )
    # Neutral geometry, not an external-paper icon.
    add_box(
        slide,
        f"{prefix}_ground",
        x + 0.10,
        y + h * 0.58,
        w - 0.20,
        h * 0.18,
        "DDE4EA",
        "DDE4EA",
        0,
        False,
    )
    triangle = set_name(
        slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(x + w * 0.42),
            Inches(y + h * 0.25),
            Inches(w * 0.16),
            Inches(h * 0.22),
        ),
        f"{prefix}_agent",
    )
    triangle.fill.solid()
    triangle.fill.fore_color.rgb = rgb(color)
    triangle.line.color.rgb = rgb(color)
    add_text(
        slide,
        f"{prefix}_label",
        x + 0.04,
        y + h - 0.26,
        w - 0.08,
        0.20,
        label,
        10.5,
        MUTED,
        True,
    )


def build_concept_slide(slide) -> None:
    """Three-region problem → method capability → evidence overview."""

    remove_all_shapes(slide)
    set_slide_white(slide)

    y, h = 0.18, 7.10
    gap = 0.10
    x_a, w_a = 0.18, 4.46
    x_b, w_b = x_a + w_a + gap, 5.26
    x_c, w_c = x_b + w_b + gap, 3.13

    add_panel(
        slide,
        "A",
        x_a,
        y,
        w_a,
        h,
        "(a) World-model logic in EO",
        BLUE_LIGHT,
        BLUE,
    )
    add_panel(
        slide,
        "B",
        x_b,
        y,
        w_b,
        h,
        "(b) TerraState exposes a testable state path",
        TEAL_LIGHT,
        TEAL,
    )
    add_panel(
        slide,
        "C",
        x_c,
        y,
        w_c,
        h,
        "(c) Hierarchical evidence",
        ORANGE_LIGHT,
        ORANGE,
    )

    # A: aligned action-conditioned and EO world-model rows.
    add_text(
        slide,
        "A_action_title",
        x_a + 0.20,
        0.82,
        w_a - 0.40,
        0.32,
        "Typical action-conditioned world model",
        15,
        BLUE,
        True,
        PP_ALIGN.LEFT,
    )
    add_scene_placeholder(
        slide,
        "A_action_scene",
        x_a + 0.22,
        1.28,
        1.00,
        1.15,
        "scene",
        BLUE,
    )
    add_state_grid(
        slide,
        "A_action_state",
        x_a + 1.63,
        1.30,
        0.86,
        1.10,
        "latent state",
        TEAL,
    )
    add_scene_placeholder(
        slide,
        "A_action_future",
        x_a + 3.04,
        1.28,
        1.00,
        1.15,
        "rollout",
        BLUE,
    )
    add_line(
        slide,
        "A_action_scene_to_state",
        x_a + 1.22,
        1.84,
        x_a + 1.63,
        1.84,
        INK,
        1.4,
    )
    add_line(
        slide,
        "A_action_state_to_future",
        x_a + 2.49,
        1.84,
        x_a + 3.04,
        1.84,
        INK,
        1.5,
    )
    add_text(
        slide,
        "A_action_condition",
        x_a + 2.43,
        2.22,
        0.78,
        0.26,
        "action a_t ↑",
        11,
        BLUE,
        True,
    )

    add_line(
        slide,
        "A_row_separator",
        x_a + 0.20,
        2.72,
        x_a + w_a - 0.20,
        2.72,
        "D0D5DD",
        0.8,
        False,
    )
    add_text(
        slide,
        "A_eo_title",
        x_a + 0.20,
        2.91,
        w_a - 0.40,
        0.34,
        "EO world modeling under exogenous forcing",
        15,
        TEAL,
        True,
        PP_ALIGN.LEFT,
    )
    add_image_slot(
        slide,
        "A_eo_history",
        x_a + 0.22,
        3.43,
        1.00,
        1.28,
        "sparse EO history",
        "project-image slot",
        BLUE,
        frames=2,
    )
    add_state_grid(
        slide,
        "A_eo_state",
        x_a + 1.63,
        3.50,
        0.86,
        1.12,
        "Earth state",
        TEAL,
    )
    add_image_slot(
        slide,
        "A_eo_future",
        x_a + 3.04,
        3.43,
        1.00,
        1.28,
        "future EO",
        "project-image slot",
        ORANGE,
        frames=1,
    )
    add_line(
        slide,
        "A_eo_history_to_state",
        x_a + 1.22,
        4.05,
        x_a + 1.63,
        4.05,
        INK,
        1.4,
    )
    add_line(
        slide,
        "A_eo_state_to_future",
        x_a + 2.49,
        4.05,
        x_a + 3.04,
        4.05,
        INK,
        1.5,
    )
    add_text(
        slide,
        "A_eo_condition",
        x_a + 2.38,
        4.43,
        0.92,
        0.28,
        "future weather ↑",
        11,
        PURPLE,
        True,
    )
    add_box(
        slide,
        "A_problem_callout",
        x_a + 0.26,
        5.18,
        w_a - 0.52,
        1.22,
        WHITE,
        RED,
        1.2,
        True,
        text=(
            "Endpoint accuracy observes the future output,\n"
            "but leaves state use and forcing use untested."
        ),
        text_size=14,
        text_color=RED,
        bold=True,
    )

    # B: compact TerraState state path and exact intervention points.
    add_text(
        slide,
        "B_summary",
        x_b + 0.28,
        0.78,
        w_b - 0.56,
        0.65,
        (
            "TerraState turns the internal pathway into a falsifiable property:\n"
            "the state must contribute, and the forecast must respond to supplied weather."
        ),
        13.5,
        INK,
        True,
    )

    chain_y = 2.74
    add_image_slot(
        slide,
        "B_history",
        x_b + 0.18,
        chain_y - 0.43,
        0.70,
        1.08,
        "history",
        "EO slot",
        BLUE,
        frames=2,
    )
    add_state_grid(
        slide,
        "B_state_zt",
        x_b + 1.10,
        chain_y - 0.38,
        0.62,
        0.98,
        "z_t",
        TEAL,
    )
    add_box(
        slide,
        "B_transition",
        x_b + 1.94,
        chain_y - 0.36,
        0.92,
        0.94,
        PURPLE_LIGHT,
        PURPLE,
        1.5,
        True,
        text="shared T\nweather-conditioned",
        text_size=12.5,
        bold=True,
    )
    add_state_grid(
        slide,
        "B_state_future",
        x_b + 3.08,
        chain_y - 0.38,
        0.62,
        0.98,
        "z_(t+h)",
        TEAL,
        True,
    )
    add_box(
        slide,
        "B_readout",
        x_b + 3.92,
        chain_y - 0.30,
        0.52,
        0.82,
        GREEN_LIGHT,
        GREEN,
        1.2,
        True,
        text="O",
        text_size=18,
        bold=True,
    )
    add_box(
        slide,
        "B_state_contribution",
        x_b + 4.61,
        chain_y - 0.30,
        0.46,
        0.82,
        ORANGE_LIGHT,
        ORANGE,
        1.2,
        True,
        text="r_h",
        text_size=15,
        bold=True,
    )
    add_line(slide, "B_h_to_z", x_b + 0.88, chain_y, x_b + 1.10, chain_y, INK, 1.4)
    add_line(slide, "B_z_to_T", x_b + 1.72, chain_y, x_b + 1.94, chain_y, INK, 1.4)
    add_line(slide, "B_T_to_zf", x_b + 2.86, chain_y, x_b + 3.08, chain_y, INK, 1.4)
    add_line(slide, "B_zf_to_O", x_b + 3.70, chain_y, x_b + 3.92, chain_y, INK, 1.4)
    add_line(slide, "B_O_to_r", x_b + 4.44, chain_y, x_b + 4.61, chain_y, INK, 1.4)

    add_weather_strip(
        slide,
        "B_future_weather",
        x_b + 1.89,
        1.53,
        1.04,
        0.74,
        "future weather",
        True,
    )
    add_line(
        slide,
        "B_weather_to_T",
        x_b + 2.41,
        2.27,
        x_b + 2.41,
        chain_y - 0.36,
        PURPLE,
        1.4,
    )
    add_intervention_port(
        slide,
        "B_Q3_port",
        x_b + 0.66,
        1.61,
        1.02,
        0.56,
        "Q3",
        "replace future weather",
        PURPLE,
    )
    add_line(
        slide,
        "B_Q3_pointer",
        x_b + 1.68,
        1.88,
        x_b + 1.89,
        1.88,
        PURPLE,
        1.0,
        False,
        MSO_LINE_DASH_STYLE.DASH,
    )

    add_plus(slide, "B_plus", x_b + 4.40, 4.10, 0.42)
    add_box(
        slide,
        "B_context_forecast",
        x_b + 1.05,
        4.15,
        1.60,
        0.64,
        BLUE_LIGHT,
        BLUE,
        1.1,
        True,
        text="context-only forecast b_h",
        text_size=11.5,
        bold=True,
    )
    add_image_slot(
        slide,
        "B_final_forecast",
        x_b + 4.82,
        3.78,
        0.36,
        1.04,
        "ŷ",
        "output",
        ORANGE,
        frames=1,
    )
    add_elbow(
        slide,
        "B_r_to_plus",
        [
            (x_b + 4.84, chain_y + 0.52),
            (x_b + 4.84, 4.31),
            (x_b + 4.82, 4.31),
        ],
        ORANGE,
        1.3,
    )
    add_line(
        slide,
        "B_b_to_plus",
        x_b + 2.65,
        4.47,
        x_b + 4.40,
        4.31,
        BLUE,
        1.2,
    )
    add_line(
        slide,
        "B_plus_to_output",
        x_b + 4.82,
        4.31,
        x_b + 4.84,
        4.31,
        ORANGE,
        1.5,
    )
    add_intervention_port(
        slide,
        "B_Q2_port",
        x_b + 2.92,
        3.64,
        1.18,
        0.62,
        "Q2",
        "remove state contribution",
        ORANGE,
    )
    add_line(
        slide,
        "B_Q2_pointer",
        x_b + 4.10,
        3.95,
        x_b + 4.72,
        3.95,
        ORANGE,
        1.0,
        False,
        MSO_LINE_DASH_STYLE.DASH,
    )
    add_box(
        slide,
        "B_capability_callout",
        x_b + 0.40,
        5.35,
        w_b - 0.80,
        0.98,
        WHITE,
        TEAL,
        1.2,
        True,
        text="Explicit state path + targeted interventions\n→ internally testable EO world model",
        text_size=14,
        text_color=TEAL,
        bold=True,
    )

    # C: evidence ladder, no Q4.
    cards = [
        (
            "Q1",
            "Forecast utility",
            "prerequisite",
            BLUE_LIGHT,
            BLUE,
        ),
        (
            "Q2",
            "Load-bearing state",
            "defining core",
            ORANGE_LIGHT,
            ORANGE,
        ),
        (
            "Q3",
            "Weather-response fidelity",
            "forcing grounding",
            PURPLE_LIGHT,
            PURPLE,
        ),
    ]
    base_y = 1.12
    for idx, (q, title, role, fill, accent) in enumerate(cards):
        yy = base_y + idx * 1.62
        add_box(
            slide,
            f"C_{q}_card",
            x_c + 0.30,
            yy,
            w_c - 0.60,
            1.18,
            fill,
            accent,
            1.4 if q == "Q2" else 1.1,
            True,
        )
        add_box(
            slide,
            f"C_{q}_badge",
            x_c + 0.45,
            yy + 0.28,
            0.52,
            0.52,
            accent,
            accent,
            0,
            True,
            text=q,
            text_size=14,
            text_color=WHITE,
            bold=True,
        )
        add_text(
            slide,
            f"C_{q}_title",
            x_c + 1.10,
            yy + 0.18,
            w_c - 1.55,
            0.44,
            title,
            15,
            INK,
            True,
            PP_ALIGN.LEFT,
        )
        add_text(
            slide,
            f"C_{q}_role",
            x_c + 1.10,
            yy + 0.62,
            w_c - 1.55,
            0.30,
            role,
            11.5,
            accent,
            True,
            PP_ALIGN.LEFT,
        )
        if idx < len(cards) - 1:
            add_line(
                slide,
                f"C_{q}_down",
                x_c + w_c / 2,
                yy + 1.18,
                x_c + w_c / 2,
                yy + 1.52,
                LINE,
                1.1,
            )
    add_text(
        slide,
        "C_evidence_note",
        x_c + 0.30,
        6.05,
        w_c - 0.60,
        0.74,
        "One selected model;\nthree increasingly specific questions.",
        12.5,
        MUTED,
        True,
    )

    # Thin cross-panel arrows make the narrative explicit.
    add_line(
        slide,
        "A_to_B_narrative",
        x_a + w_a + 0.01,
        6.76,
        x_b - 0.01,
        6.76,
        TEAL,
        1.6,
    )
    add_line(
        slide,
        "B_to_C_narrative",
        x_b + w_b + 0.01,
        6.76,
        x_c - 0.01,
        6.76,
        ORANGE,
        1.6,
    )


def add_compact_panel(slide, prefix, x, y, w, h, title, tint, accent):
    add_box(slide, f"{prefix}_panel", x, y, w, h, "FCFDFE", "D0D5DD", 0.9, True)
    add_box(
        slide,
        f"{prefix}_header",
        x,
        y,
        w,
        0.43,
        tint,
        accent,
        1.0,
        True,
        text=title,
        text_size=14.5,
        text_color=INK,
        bold=True,
    )


def add_module(slide, name, x, y, w, h, label, accent=TEAL, tint=TEAL_LIGHT, size=13):
    return add_box(
        slide, name, x, y, w, h, tint, accent, 1.3, True,
        text=label, text_size=size, text_color=INK, bold=True,
    )


def build_method_slide_v2(slide) -> None:
    """Compact continuous method flow inspired by strong conference figures."""

    remove_all_shapes(slide)
    set_slide_white(slide)

    y, h, gap = 0.16, 7.14, 0.08
    x_a, w_a = 0.16, 2.32
    x_b, w_b = x_a + w_a + gap, 3.08
    x_c, w_c = x_b + w_b + gap, 4.32
    x_d, w_d = x_c + w_c + gap, 2.99
    add_compact_panel(slide, "A", x_a, y, w_a, h, "(a) Context", BLUE_LIGHT, BLUE)
    add_compact_panel(slide, "B", x_b, y, w_b, h, "(b) State construction", TEAL_LIGHT, TEAL)
    add_compact_panel(slide, "C", x_c, y, w_c, h, "(c) Weather-conditioned dynamics", PURPLE_LIGHT, PURPLE)
    add_compact_panel(slide, "D", x_d, y, w_d, h, "(d) Forecast readout", ORANGE_LIGHT, ORANGE)

    # A — honest, replaceable project-image anchors.
    add_image_slot(slide, "A_history", x_a + 0.20, 0.78, 1.92, 1.28, "EO history", "project frames", BLUE, frames=3)
    add_image_slot(slide, "A_environment", x_a + 0.20, 2.24, 1.92, 1.05, "context history", "mask + past weather", BLUE, frames=2)
    add_small_map_pair(slide, "A_geography", x_a + 0.20, 3.49, 1.92, 1.02, "geography g: land cover · DEM")
    add_weather_strip(slide, "A_future_weather", x_a + 0.20, 4.72, 1.92, 1.10, "future forcing u(t:t+h)", compact=True)
    add_text(slide, "A_note", x_a + 0.24, 6.12, 1.84, 0.48, "Observed history\n+ known exogenous drivers", 11.5, MUTED, True)
    add_line(slide, "A_to_B", x_a + w_a - 0.01, 3.48, x_b + 0.18, 3.48, BLUE, 1.8)

    # B — one encoder, two historical-only branches.
    add_module(slide, "B_q", x_b + 0.28, 1.05, 2.52, 0.92, "History encoder q\ncloud-aware context", BLUE, BLUE_LIGHT, 15)
    add_token_wall(slide, "B_features", x_b + 0.48, 2.27, cols=4, rows=3, cell=0.14, gap=0.04)
    add_text(slide, "B_features_label", x_b + 0.34, 2.82, 1.02, 0.38, "context\nfeatures", 11, MUTED, False)
    add_line(slide, "B_q_to_features", x_b + 1.54, 1.97, x_b + 0.82, 2.27, BLUE, 1.3)

    add_text(slide, "B_state_branch", x_b + 1.25, 2.30, 1.52, 0.35, "state branch", 11.5, TEAL, True)
    add_module(slide, "B_P", x_b + 1.36, 2.78, 0.72, 0.78, "P", TEAL, TEAL_LIGHT, 18)
    add_state_grid(slide, "B_zt", x_b + 2.24, 2.58, 0.66, 1.22, "zₜ", TEAL, False)
    add_line(slide, "B_features_to_P", x_b + 1.10, 2.54, x_b + 1.36, 3.17, TEAL, 1.4)
    add_line(slide, "B_P_to_zt", x_b + 2.08, 3.17, x_b + 2.24, 3.17, TEAL, 1.5)

    add_text(slide, "B_base_branch", x_b + 0.27, 4.18, 2.54, 0.35, "context-only reference branch", 11.5, BLUE, True)
    add_module(slide, "B_baseline", x_b + 0.46, 4.65, 2.18, 0.86, "bₕ\ncontext-only forecast", BLUE, BLUE_LIGHT, 14)
    add_elbow(slide, "B_features_to_base", [(x_b + 0.82, 2.71), (x_b + 0.82, 4.33), (x_b + 1.55, 4.65)], BLUE, 1.25)
    add_text(slide, "B_note", x_b + 0.35, 5.94, 2.38, 0.62, "q and b_h read historical\ninformation only", 11.5, MUTED, False)
    add_line(slide, "B_zt_to_C", x_b + 2.90, 3.17, x_c + 0.25, 3.17, TEAL, 2.0)

    # C — conditioning row converges on one shared transition.
    add_weather_strip(slide, "C_weather", x_c + 0.30, 0.86, 1.70, 1.08, "future weather u", compact=True)
    add_small_map_pair(slide, "C_geo", x_c + 2.15, 0.86, 1.22, 1.08, "geography g")
    add_clock_chip(slide, "C_h", x_c + 3.53, 0.86, 0.52, 1.08, "h")
    add_intervention_port(slide, "C_Q3", x_c + 0.35, 2.13, 1.52, 0.62, "Q3", "replace future weather", PURPLE)
    add_module(slide, "C_T", x_c + 1.15, 2.78, 1.90, 1.62, "Shared transition T\nstate × forcing", PURPLE, PURPLE_LIGHT, 17)
    add_line(slide, "C_weather_to_T", x_c + 1.15, 1.94, x_c + 1.72, 2.78, ORANGE, 1.4)
    add_line(slide, "C_geo_to_T", x_c + 2.76, 1.94, x_c + 2.56, 2.78, GREEN, 1.35)
    add_line(slide, "C_h_to_T", x_c + 3.79, 1.94, x_c + 2.99, 2.94, MUTED, 1.25)
    add_line(slide, "C_Q3_to_weather", x_c + 1.10, 2.13, x_c + 1.10, 1.94, PURPLE, 1.0, True, MSO_LINE_DASH_STYLE.DASH)
    add_state_grid(slide, "C_zfuture", x_c + 3.32, 2.87, 0.70, 1.42, "zₜ₊ₕ", TEAL, True)
    add_line(slide, "C_T_to_zfuture", x_c + 3.05, 3.59, x_c + 3.32, 3.59, TEAL, 1.8)
    add_text(slide, "C_shared_note", x_c + 0.38, 4.91, 3.55, 0.62, "One weather-conditioned transition\nis reused for every queried horizon.", 12.5, PURPLE, True)
    add_text(slide, "C_condition_note", x_c + 0.48, 5.85, 3.35, 0.48, "Explicit inputs: u(t:t+h), g, and h", 11.5, MUTED, False)
    add_line(slide, "C_to_D", x_c + w_c - 0.01, 3.59, x_d + 0.16, 3.59, TEAL, 1.8)

    # D — read out state contribution, fuse with context-only forecast.
    add_module(slide, "D_O", x_d + 0.23, 2.78, 0.64, 1.08, "O", GREEN, GREEN_LIGHT, 20)
    add_module(slide, "D_r", x_d + 1.02, 2.78, 0.72, 1.08, "rₕ", TEAL, TEAL_LIGHT, 17)
    add_line(slide, "D_O_to_r", x_d + 0.87, 3.32, x_d + 1.02, 3.32, GREEN, 1.5)
    add_plus(slide, "D_plus", x_d + 1.92, 3.09, 0.48)
    add_line(slide, "D_r_to_plus", x_d + 1.74, 3.32, x_d + 1.92, 3.32, ORANGE, 1.5)
    add_intervention_port(slide, "D_Q2", x_d + 0.72, 1.28, 1.52, 0.70, "Q2", "remove state contribution", ORANGE)
    add_line(slide, "D_Q2_to_r", x_d + 1.48, 1.98, x_d + 1.38, 2.78, ORANGE, 1.0, True, MSO_LINE_DASH_STYLE.DASH)
    add_elbow(slide, "D_baseline_path", [(x_b + 2.64, 5.08), (x_b + 2.78, 4.72), (x_d + 2.16, 4.72), (x_d + 2.16, 3.57)], BLUE, 1.35)
    add_text(slide, "D_equation", x_d + 0.35, 4.12, 2.30, 0.48, "ŷ(t+h) = bₕ + rₕ", 17, ORANGE, True)
    for idx, horizon in enumerate(("h=5", "h=10", "h=20")):
        add_image_slot(slide, f"D_output_{idx}", x_d + 0.20 + idx * 0.92, 5.02, 0.78, 1.08, horizon, "NDVI", ORANGE)
    add_text(slide, "D_output_label", x_d + 0.38, 6.28, 2.22, 0.34, "Vegetation forecasts", 13, ORANGE, True)


def build_concept_slide_v2(slide) -> None:
    """Concise problem → explicit state path → evidence definition figure."""

    remove_all_shapes(slide)
    set_slide_white(slide)
    y, h, gap = 0.16, 7.14, 0.08
    x_a, w_a = 0.16, 4.38
    x_b, w_b = x_a + w_a + gap, 5.52
    x_c, w_c = x_b + w_b + gap, 3.03
    add_compact_panel(slide, "A", x_a, y, w_a, h, "(a) Prediction vs. world modeling", BLUE_LIGHT, BLUE)
    add_compact_panel(slide, "B", x_b, y, w_b, h, "(b) TerraState: a testable state path", TEAL_LIGHT, TEAL)
    add_compact_panel(slide, "C", x_c, y, w_c, h, "(c) Evidence", ORANGE_LIGHT, ORANGE)

    # A1: conventional endpoint predictor.
    add_text(slide, "A1_title", x_a + 0.24, 0.80, 3.90, 0.30, "Conventional fixed-horizon EO predictor", 14, BLUE, True, PP_ALIGN.LEFT)
    add_image_slot(slide, "A1_history", x_a + 0.24, 1.26, 1.04, 1.06, "history", "EO frames", BLUE, frames=2)
    add_module(slide, "A1_F", x_a + 1.62, 1.42, 1.02, 0.74, "Fₕ", BLUE, BLUE_LIGHT, 17)
    add_image_slot(slide, "A1_output", x_a + 2.99, 1.26, 1.04, 1.06, "forecast", "future EO", BLUE)
    add_line(slide, "A1_h_to_F", x_a + 1.28, 1.79, x_a + 1.62, 1.79, BLUE, 1.5)
    add_line(slide, "A1_F_to_out", x_a + 2.64, 1.79, x_a + 2.99, 1.79, BLUE, 1.5)
    add_box(slide, "A1_callout", x_a + 0.40, 2.56, 3.56, 0.62, "FFF6F5", RED, 1.0, True, text="Only the endpoint error is observed", text_size=12.5, text_color=RED, bold=True)

    add_line(slide, "A_divider", x_a + 0.24, 3.48, x_a + 4.14, 3.48, "D0D5DD", 0.9, False)

    # A2: EO world-model semantics.
    add_text(slide, "A2_title", x_a + 0.24, 3.73, 3.90, 0.30, "EO world model under exogenous forcing", 14, TEAL, True, PP_ALIGN.LEFT)
    add_image_slot(slide, "A2_history", x_a + 0.24, 4.32, 0.88, 1.02, "history", "EO", BLUE, frames=2)
    add_state_grid(slide, "A2_zt", x_a + 1.35, 4.26, 0.68, 1.14, "zₜ", TEAL, False)
    add_module(slide, "A2_T", x_a + 2.28, 4.30, 0.80, 1.04, "T", PURPLE, PURPLE_LIGHT, 19)
    add_state_grid(slide, "A2_zf", x_a + 3.34, 4.26, 0.68, 1.14, "zₜ₊ₕ", TEAL, True)
    add_line(slide, "A2_hist_to_z", x_a + 1.12, 4.83, x_a + 1.35, 4.83, TEAL, 1.4)
    add_line(slide, "A2_z_to_T", x_a + 2.03, 4.83, x_a + 2.28, 4.83, TEAL, 1.4)
    add_line(slide, "A2_T_to_zf", x_a + 3.08, 4.83, x_a + 3.34, 4.83, TEAL, 1.4)
    add_weather_strip(slide, "A2_weather", x_a + 1.88, 5.70, 1.54, 0.78, "future weather u", compact=True)
    add_line(slide, "A2_weather_to_T", x_a + 2.65, 5.70, x_a + 2.68, 5.34, PURPLE, 1.2)
    add_box(slide, "A2_callout", x_a + 0.42, 6.62, 3.52, 0.44, TEAL_LIGHT, TEAL, 1.0, True, text="State use and forcing use become observable", text_size=11.5, text_color=TEAL, bold=True)

    # B: the concrete TerraState chain; compact labels and intervention ports.
    add_text(slide, "B_subtitle", x_b + 0.30, 0.79, w_b - 0.60, 0.56, "A predictive state should carry forecast information\nand respond to the supplied future weather.", 13, INK, True)
    chain_y = 2.58
    add_image_slot(slide, "B_history", x_b + 0.28, chain_y, 0.78, 1.16, "history", "EO", BLUE, frames=2)
    add_module(slide, "B_qP", x_b + 1.27, chain_y + 0.13, 0.78, 0.90, "q → P", BLUE, BLUE_LIGHT, 14)
    add_state_grid(slide, "B_zt", x_b + 2.25, chain_y, 0.68, 1.16, "zₜ", TEAL, False)
    add_module(slide, "B_T", x_b + 3.15, chain_y - 0.06, 0.90, 1.28, "shared\nT", PURPLE, PURPLE_LIGHT, 15)
    add_state_grid(slide, "B_zf", x_b + 4.27, chain_y, 0.68, 1.16, "zₜ₊ₕ", TEAL, True)
    add_module(slide, "B_O", x_b + 5.08, chain_y + 0.13, 0.30, 0.90, "O", GREEN, GREEN_LIGHT, 14)
    for idx, (a, b) in enumerate(((1.06, 1.27), (2.05, 2.25), (2.93, 3.15), (4.05, 4.27), (4.95, 5.08)), 1):
        add_line(slide, f"B_chain_{idx}", x_b + a, chain_y + 0.58, x_b + b, chain_y + 0.58, TEAL, 1.45)

    add_weather_strip(slide, "B_weather", x_b + 2.95, 1.48, 1.38, 0.78, "future weather u", compact=True)
    add_line(slide, "B_weather_to_T", x_b + 3.64, 2.26, x_b + 3.60, chain_y - 0.06, PURPLE, 1.25)
    add_intervention_port(slide, "B_Q3", x_b + 1.64, 1.54, 1.08, 0.64, "Q3", "replace future weather", PURPLE)
    add_line(slide, "B_Q3_pointer", x_b + 2.72, 1.86, x_b + 2.95, 1.86, PURPLE, 1.0, True, MSO_LINE_DASH_STYLE.DASH)

    add_module(slide, "B_base", x_b + 0.72, 4.48, 1.86, 0.72, "context-only forecast bₕ", BLUE, BLUE_LIGHT, 12.5)
    add_plus(slide, "B_plus", x_b + 4.47, 4.59, 0.48)
    add_module(slide, "B_output", x_b + 5.02, 4.50, 0.40, 0.66, "ŷ", ORANGE, ORANGE_LIGHT, 16)
    add_elbow(slide, "B_state_to_plus", [(x_b + 5.23, chain_y + 1.03), (x_b + 5.23, 4.83), (x_b + 4.95, 4.83)], ORANGE, 1.35)
    add_text(slide, "B_state_contribution", x_b + 4.78, 4.00, 0.42, 0.28, "rₕ", 12.5, ORANGE, True)
    add_line(slide, "B_base_to_plus", x_b + 2.58, 4.84, x_b + 4.47, 4.84, BLUE, 1.25)
    add_line(slide, "B_plus_to_output", x_b + 4.95, 4.83, x_b + 5.05, 4.83, ORANGE, 1.4)
    add_intervention_port(slide, "B_Q2", x_b + 2.72, 5.45, 1.44, 0.66, "Q2", "remove state contribution", ORANGE)
    add_line(slide, "B_Q2_pointer", x_b + 4.16, 5.78, x_b + 4.71, 5.05, ORANGE, 1.0, True, MSO_LINE_DASH_STYLE.DASH)
    add_box(slide, "B_claim", x_b + 0.58, 6.40, 4.38, 0.56, TEAL_LIGHT, TEAL, 1.1, True, text="Explicit path + targeted interventions → internally testable", text_size=12.5, text_color=TEAL, bold=True)

    # C: three questions in increasing specificity, no Q4.
    cards = [
        ("Q1", "Forecast utility", "prerequisite", BLUE, BLUE_LIGHT),
        ("Q2", "Load-bearing state", "defining core", ORANGE, ORANGE_LIGHT),
        ("Q3", "Weather response", "forcing fidelity", PURPLE, PURPLE_LIGHT),
    ]
    for idx, (q, title, role, accent, tint) in enumerate(cards):
        yy = 1.05 + idx * 1.72
        add_box(slide, f"C_{q}", x_c + 0.26, yy, 2.51, 1.26, tint, accent, 1.3, True)
        add_box(slide, f"C_{q}_badge", x_c + 0.42, yy + 0.29, 0.58, 0.68, accent, accent, 0, True, text=q, text_size=15, text_color=WHITE, bold=True)
        add_text(slide, f"C_{q}_title", x_c + 1.13, yy + 0.18, 1.46, 0.48, title, 13.5, INK, True, PP_ALIGN.LEFT)
        add_text(slide, f"C_{q}_role", x_c + 1.13, yy + 0.69, 1.46, 0.30, role, 11.5, accent, True, PP_ALIGN.LEFT)
        if idx < 2:
            add_line(slide, f"C_flow_{idx}", x_c + 1.52, yy + 1.27, x_c + 1.52, yy + 1.66, "98A2B3", 1.0)
    add_text(slide, "C_note", x_c + 0.40, 6.45, 2.25, 0.50, "One selected model;\nthree increasingly specific tests.", 11.5, MUTED, True)


def build_concept_slide_thesis(slide) -> None:
    """Alternative: thesis-first upper row plus three visual evidence cards."""

    remove_all_shapes(slide)
    set_slide_white(slide)
    gap = 0.10
    y_top, h_top = 0.18, 3.02
    x_a, w_a = 0.18, 3.42
    x_b, w_b = x_a + w_a + gap, 5.48
    x_c, w_c = x_b + w_b + gap, 3.97
    add_compact_panel(slide, "TA", x_a, y_top, w_a, h_top, "(a) Endpoint prediction", BLUE_LIGHT, BLUE)
    add_compact_panel(slide, "TB", x_b, y_top, w_b, h_top, "(b) TerraState world model", TEAL_LIGHT, TEAL)
    add_compact_panel(slide, "TC", x_c, y_top, w_c, h_top, "(c) Testable properties", ORANGE_LIGHT, ORANGE)

    add_image_slot(slide, "TA_hist", x_a + 0.22, 0.92, 0.90, 1.05, "history", "EO", BLUE, frames=2)
    add_module(slide, "TA_F", x_a + 1.37, 1.04, 0.70, 0.78, "Fₕ", BLUE, BLUE_LIGHT, 17)
    add_image_slot(slide, "TA_out", x_a + 2.32, 0.92, 0.88, 1.05, "future", "EO", BLUE)
    add_line(slide, "TA_1", x_a + 1.12, 1.45, x_a + 1.37, 1.45, BLUE, 1.4)
    add_line(slide, "TA_2", x_a + 2.07, 1.45, x_a + 2.32, 1.45, BLUE, 1.4)
    add_box(slide, "TA_note", x_a + 0.36, 2.25, 2.70, 0.50, "FFF6F5", RED, 1.0, True, text="Endpoint score only", text_size=13, text_color=RED, bold=True)

    chain_y = 1.31
    nodes = [
        ("TB_hist", 0.22, 0.76, "history", BLUE, BLUE_LIGHT),
        ("TB_qP", 1.13, 0.76, "q→P", BLUE, BLUE_LIGHT),
        ("TB_zt", 2.04, 0.72, "zₜ", TEAL, TEAL_LIGHT),
        ("TB_T", 2.91, 0.82, "T", PURPLE, PURPLE_LIGHT),
        ("TB_zf", 3.88, 0.72, "zₜ₊ₕ", TEAL, TEAL_LIGHT),
        ("TB_O", 4.79, 0.44, "O", GREEN, GREEN_LIGHT),
    ]
    for name, dx, ww, label, accent, tint in nodes:
        add_module(slide, name, x_b + dx, chain_y, ww, 0.86, label, accent, tint, 14.5)
    for idx, (x1, x2) in enumerate(((0.98, 1.13), (1.89, 2.04), (2.76, 2.91), (3.73, 3.88), (4.60, 4.79)), 1):
        add_line(slide, f"TB_flow_{idx}", x_b + x1, chain_y + 0.43, x_b + x2, chain_y + 0.43, TEAL, 1.35)
    add_weather_strip(slide, "TB_weather", x_b + 2.62, 0.55, 1.38, 0.62, "future weather u", compact=True)
    add_line(slide, "TB_weather_T", x_b + 3.31, 1.17, x_b + 3.32, chain_y, PURPLE, 1.1)
    add_box(slide, "TB_claim", x_b + 0.56, 2.39, 4.34, 0.42, TEAL_LIGHT, TEAL, 1.0, True, text="Explicit state path under exogenous forcing", text_size=12.5, text_color=TEAL, bold=True)

    criteria = [
        ("1", "State carries\nforecast information", TEAL),
        ("2", "Weather changes\nthe predicted response", PURPLE),
    ]
    for idx, (num, text_value, accent) in enumerate(criteria):
        yy = 0.86 + idx * 1.02
        add_box(slide, f"TC_card_{idx}", x_c + 0.30, yy, 3.37, 0.80, WHITE, accent, 1.1, True)
        add_box(slide, f"TC_num_{idx}", x_c + 0.46, yy + 0.17, 0.44, 0.44, accent, accent, 0, True, text=num, text_size=13, text_color=WHITE, bold=True)
        add_text(slide, f"TC_text_{idx}", x_c + 1.08, yy + 0.10, 2.28, 0.60, text_value, 12, INK, True, PP_ALIGN.LEFT)

    # Bottom: visual evidence cards, each with a miniature intervention diagram.
    add_text(slide, "TE_title", 0.30, 3.43, 12.72, 0.34, "One selected model, three increasingly specific questions", 14, INK, True)
    card_y, card_h, card_w = 3.92, 3.18, 4.13
    card_specs = [
        ("Q1", "Forecast utility", "Does the complete model retain useful EO forecasting skill?", BLUE, BLUE_LIGHT),
        ("Q2", "Load-bearing state", "Does removing the state contribution degrade the forecast?", ORANGE, ORANGE_LIGHT),
        ("Q3", "Weather-response fidelity", "Does replacing future weather alter the response correctly?", PURPLE, PURPLE_LIGHT),
    ]
    for idx, (q, title, question, accent, tint) in enumerate(card_specs):
        xx = 0.24 + idx * (card_w + 0.16)
        add_box(slide, f"TE_{q}", xx, card_y, card_w, card_h, tint, accent, 1.2, True)
        add_box(slide, f"TE_{q}_badge", xx + 0.18, card_y + 0.18, 0.54, 0.54, accent, accent, 0, True, text=q, text_size=13, text_color=WHITE, bold=True)
        add_text(slide, f"TE_{q}_title", xx + 0.86, card_y + 0.12, 2.98, 0.42, title, 13.5, INK, True, PP_ALIGN.LEFT)
        add_text(slide, f"TE_{q}_question", xx + 0.28, card_y + 0.82, 3.57, 0.64, question, 11.5, MUTED, False)
        if q == "Q1":
            add_image_slot(slide, "TE_Q1_pred", xx + 0.52, card_y + 1.72, 1.03, 0.98, "forecast", "NDVI", BLUE)
            add_image_slot(slide, "TE_Q1_target", xx + 2.45, card_y + 1.72, 1.03, 0.98, "target", "NDVI", TEAL)
            add_line(slide, "TE_Q1_compare", xx + 1.55, card_y + 2.21, xx + 2.45, card_y + 2.21, BLUE, 1.2)
        elif q == "Q2":
            add_module(slide, "TE_Q2_full", xx + 0.46, card_y + 1.75, 1.20, 0.86, "full\nstate path", TEAL, TEAL_LIGHT, 12)
            add_module(slide, "TE_Q2_cut", xx + 2.50, card_y + 1.75, 1.20, 0.86, "state\nremoved", ORANGE, WHITE, 12)
            add_text(slide, "TE_Q2_delta", xx + 1.73, card_y + 1.92, 0.70, 0.42, "Δ", 18, ORANGE, True)
        else:
            add_weather_strip(slide, "TE_Q3_real", xx + 0.34, card_y + 1.65, 1.42, 0.90, "actual u", compact=True)
            add_weather_strip(slide, "TE_Q3_donor", xx + 2.35, card_y + 1.65, 1.42, 0.90, "replaced u′", compact=True)
            add_text(slide, "TE_Q3_delta", xx + 1.76, card_y + 1.90, 0.60, 0.36, "→ Δŷ", 12.5, PURPLE, True)


def build_method_slide_dense(slide) -> None:
    """Alternative: one continuous AAAI-style architecture pipeline."""

    remove_all_shapes(slide)
    set_slide_white(slide)
    # Soft stage backplates and ribbon headers.
    stages = [
        (0.15, 2.28, "(a) Inputs", BLUE_LIGHT, BLUE),
        (2.51, 3.76, "(b) State construction", TEAL_LIGHT, TEAL),
        (6.35, 3.74, "(c) Shared dynamics", PURPLE_LIGHT, PURPLE),
        (10.17, 3.01, "(d) Readout", ORANGE_LIGHT, ORANGE),
    ]
    for idx, (xx, ww, title, tint, accent) in enumerate(stages):
        add_box(slide, f"MD_stage_{idx}", xx, 0.18, ww, 7.10, "FCFDFE", "D0D5DD", 0.8, True)
        add_box(slide, f"MD_header_{idx}", xx, 0.18, ww, 0.42, tint, accent, 1.0, True, text=title, text_size=14, text_color=INK, bold=True)

    # Input anchors.
    add_image_slot(slide, "MD_history", 0.36, 0.94, 1.76, 1.30, "EO history", "Sentinel-2 sequence", BLUE, frames=3)
    add_weather_strip(slide, "MD_weather", 0.36, 2.58, 1.76, 1.08, "future weather u", compact=True)
    add_small_map_pair(slide, "MD_geo", 0.36, 4.00, 1.76, 1.02, "geography g")
    add_clock_chip(slide, "MD_h", 0.86, 5.40, 0.74, 0.80, "h")

    # Continuous main chain.
    cy = 2.46
    add_module(slide, "MD_q", 2.76, cy, 1.05, 1.08, "history\nencoder q", BLUE, BLUE_LIGHT, 14)
    add_token_wall(slide, "MD_context_tokens", 4.05, cy + 0.15, cols=4, rows=4, cell=0.13, gap=0.035)
    add_module(slide, "MD_P", 4.83, cy + 0.14, 0.58, 0.80, "P", TEAL, TEAL_LIGHT, 17)
    add_state_grid(slide, "MD_zt", 5.56, cy - 0.02, 0.68, 1.12, "zₜ", TEAL, False)
    add_module(slide, "MD_T", 7.23, cy - 0.28, 1.45, 1.64, "Shared T\nstate × forcing", PURPLE, PURPLE_LIGHT, 16)
    add_state_grid(slide, "MD_zf", 9.18, cy - 0.02, 0.68, 1.12, "zₜ₊ₕ", TEAL, True)
    add_module(slide, "MD_O", 10.45, cy, 0.58, 1.08, "O", GREEN, GREEN_LIGHT, 19)
    add_module(slide, "MD_r", 11.28, cy, 0.66, 1.08, "rₕ", TEAL, TEAL_LIGHT, 16)
    add_plus(slide, "MD_plus", 12.16, cy + 0.30, 0.48)
    chain_pairs = [
        (3.81, 4.05, BLUE), (4.67, 4.83, TEAL),
        (5.41, 5.56, TEAL), (6.24, 7.23, TEAL), (8.68, 9.18, TEAL),
        (9.86, 10.45, TEAL), (11.03, 11.28, GREEN), (11.94, 12.16, ORANGE),
    ]
    for idx, (a, b, color) in enumerate(chain_pairs):
        add_line(slide, f"MD_chain_{idx}", a, cy + 0.54, b, cy + 0.54, color, 1.5)
    add_elbow(
        slide,
        "MD_history_to_q",
        [(2.12, 1.59), (2.30, 1.59), (2.30, 2.28), (3.28, 2.28), (3.28, cy)],
        BLUE,
        1.4,
    )

    # Conditioning convergence and compact verification interfaces.
    add_elbow(
        slide,
        "MD_weather_T",
        [(2.12, 3.04), (2.28, 3.04), (2.28, 1.86), (7.52, 1.86), (7.52, cy - 0.28)],
        ORANGE,
        1.20,
    )
    add_elbow(
        slide,
        "MD_geo_T",
        [(2.12, 4.51), (2.36, 4.51), (2.36, 1.98), (8.00, 1.98), (8.00, cy - 0.28)],
        GREEN,
        1.10,
    )
    add_elbow(
        slide,
        "MD_h_T",
        [(1.60, 5.80), (2.44, 5.80), (2.44, 2.10), (8.40, 2.10), (8.40, cy - 0.18)],
        MUTED,
        1.00,
    )
    add_intervention_port(slide, "MD_Q3", 6.52, 0.98, 1.35, 0.68, "Q3", "replace u", PURPLE)
    add_line(slide, "MD_Q3_ptr", 7.20, 1.66, 7.62, cy - 0.28, PURPLE, 1.0, True, MSO_LINE_DASH_STYLE.DASH)
    add_intervention_port(slide, "MD_Q2", 10.80, 0.98, 1.35, 0.68, "Q2", "remove rₕ", ORANGE)
    add_line(slide, "MD_Q2_ptr", 11.47, 1.66, 11.61, cy, ORANGE, 1.0, True, MSO_LINE_DASH_STYLE.DASH)
    add_text(slide, "MD_condition_labels", 7.17, 2.08, 1.45, 0.24, "u  ·  g  ·  h", 11.5, PURPLE, True)

    # Context-only branch and multi-horizon outputs.
    add_module(slide, "MD_base", 3.22, 5.14, 2.10, 0.88, "context-only forecast bₕ", BLUE, BLUE_LIGHT, 13)
    add_elbow(slide, "MD_base_from_q", [(3.29, 3.54), (3.29, 4.70), (4.27, 5.14)], BLUE, 1.2)
    add_elbow(slide, "MD_base_to_plus", [(5.32, 5.58), (5.52, 4.82), (12.40, 4.82), (12.40, cy + 0.78)], BLUE, 1.3)
    add_text(slide, "MD_eq", 10.62, 4.20, 2.22, 0.44, "ŷ(t+h) = bₕ + rₕ", 16, ORANGE, True)
    for idx, horizon in enumerate(("h=5", "h=10", "h=20")):
        add_image_slot(slide, f"MD_out_{idx}", 10.46 + idx * 0.84, 5.10, 0.70, 1.02, horizon, "NDVI", ORANGE)
    add_text(slide, "MD_output_label", 10.52, 6.34, 2.28, 0.34, "Vegetation forecasts", 13, ORANGE, True)
    add_text(slide, "MD_shared", 6.32, 4.02, 3.04, 0.58, "The same weather-conditioned transition\nserves every queried horizon.", 12, PURPLE, True)


def build_method_slide_multilevel(slide) -> None:
    """Alternative: overall inference chain plus three architectural zoom-ins."""

    remove_all_shapes(slide)
    set_slide_white(slide)

    # (a) Compact overall path.
    add_compact_panel(slide, "ML_A", 0.16, 0.15, 13.01, 3.10, "(a) Overall TerraState inference path", TEAL_LIGHT, TEAL)
    cy = 1.65
    add_image_slot(slide, "ML_hist", 0.38, 0.82, 1.22, 1.30, "EO history", "Sentinel-2", BLUE, frames=3)
    add_module(slide, "ML_qP", 1.90, cy - 0.36, 0.88, 0.82, "q → P", BLUE, BLUE_LIGHT, 14)
    add_state_grid(slide, "ML_zt", 3.07, cy - 0.52, 0.68, 1.14, "zₜ", TEAL, False)
    add_module(slide, "ML_T", 5.22, cy - 0.65, 1.30, 1.40, "Shared T", PURPLE, PURPLE_LIGHT, 16)
    add_state_grid(slide, "ML_zf", 7.08, cy - 0.52, 0.68, 1.14, "zₜ₊ₕ", TEAL, True)
    add_module(slide, "ML_O", 8.13, cy - 0.36, 0.52, 0.82, "O", GREEN, GREEN_LIGHT, 17)
    add_module(slide, "ML_r", 8.94, cy - 0.36, 0.62, 0.82, "rₕ", TEAL, TEAL_LIGHT, 15)
    add_plus(slide, "ML_plus", 9.87, cy - 0.19, 0.46)
    add_image_slot(slide, "ML_out", 10.70, 0.95, 1.18, 1.22, "ŷ(t+h)", "NDVI", ORANGE)
    pairs = [
        (1.60, 1.90, BLUE), (2.78, 3.07, TEAL), (3.75, 5.22, TEAL),
        (6.52, 7.08, TEAL), (7.76, 8.13, TEAL), (8.65, 8.94, GREEN),
        (9.56, 9.87, ORANGE), (10.33, 10.70, ORANGE),
    ]
    for idx, (a, b, color) in enumerate(pairs):
        add_line(slide, f"ML_chain_{idx}", a, cy, b, cy, color, 1.45)

    add_weather_strip(slide, "ML_u", 4.04, 0.72, 1.00, 0.72, "future u", compact=True)
    add_small_map_pair(slide, "ML_g", 4.02, 1.82, 0.98, 0.76, "geo g")
    add_clock_chip(slide, "ML_h", 5.30, 0.70, 0.48, 0.66, "h")
    add_line(slide, "ML_u_T", 5.04, 1.09, 5.48, cy - 0.65, ORANGE, 1.05)
    add_line(slide, "ML_g_T", 5.00, 2.20, 5.71, cy + 0.75, GREEN, 1.05)
    add_line(slide, "ML_h_T", 5.54, 1.36, 5.85, cy - 0.65, MUTED, 1.00)

    add_module(slide, "ML_b", 1.96, 2.40, 2.05, 0.52, "context-only forecast bₕ", BLUE, BLUE_LIGHT, 12)
    add_elbow(slide, "ML_q_b", [(2.34, 2.11), (2.34, 2.28), (2.98, 2.40)], BLUE, 1.05)
    add_elbow(slide, "ML_b_plus", [(4.01, 2.66), (10.10, 2.66), (10.10, cy + 0.27)], BLUE, 1.15)
    add_intervention_port(slide, "ML_Q3", 6.76, 0.60, 1.04, 0.56, "Q3", "replace u", PURPLE)
    add_line(slide, "ML_Q3_ptr", 7.28, 1.16, 6.15, cy - 0.65, PURPLE, 0.9, True, MSO_LINE_DASH_STYLE.DASH)
    add_intervention_port(slide, "ML_Q2", 9.00, 0.60, 1.04, 0.56, "Q2", "remove rₕ", ORANGE)
    add_line(slide, "ML_Q2_ptr", 9.52, 1.16, 9.25, cy - 0.36, ORANGE, 0.9, True, MSO_LINE_DASH_STYLE.DASH)
    add_text(slide, "ML_equation", 11.96, 1.20, 0.94, 0.86, "ŷ = bₕ\n+ rₕ", 14, ORANGE, True)

    # (b) Zoom: history-only state construction.
    yb, hb = 3.43, 3.87
    add_compact_panel(slide, "ML_B", 0.16, yb, 4.06, hb, "(b) History-only state construction", BLUE_LIGHT, BLUE)
    add_text(slide, "ML_B_note", 0.42, 4.02, 3.54, 0.44, "Cloud-aware context tokens are shared by both branches.", 11.5, MUTED, False)
    add_token_wall(slide, "ML_B_tokens", 0.52, 4.82, cols=5, rows=4, cell=0.17, gap=0.045)
    add_text(slide, "ML_B_token_label", 0.40, 5.64, 1.25, 0.34, "context tokens", 12, BLUE, True)
    add_module(slide, "ML_B_P", 1.82, 4.95, 0.72, 0.76, "P", TEAL, TEAL_LIGHT, 18)
    add_state_grid(slide, "ML_B_zt", 2.84, 4.73, 0.76, 1.20, "zₜ", TEAL, False)
    add_line(slide, "ML_B_tokens_P", 1.48, 5.33, 1.82, 5.33, TEAL, 1.3)
    add_line(slide, "ML_B_P_zt", 2.54, 5.33, 2.84, 5.33, TEAL, 1.3)
    add_module(slide, "ML_B_base", 1.18, 6.33, 1.86, 0.52, "context forecast bₕ", BLUE, BLUE_LIGHT, 12)
    add_elbow(slide, "ML_B_tokens_base", [(1.06, 5.68), (1.06, 6.12), (2.11, 6.33)], BLUE, 1.05)

    # (c) Zoom: condition encoding and residual shared transition.
    add_compact_panel(slide, "ML_C", 4.31, yb, 4.56, hb, "(c) Shared weather-conditioned transition", PURPLE_LIGHT, PURPLE)
    add_weather_strip(slide, "ML_C_u", 4.58, 4.10, 1.20, 0.84, "ordered u", compact=True)
    add_small_map_pair(slide, "ML_C_g", 4.58, 5.25, 1.20, 0.82, "geography g")
    add_clock_chip(slide, "ML_C_h", 4.88, 6.38, 0.62, 0.60, "h")
    add_module(slide, "ML_C_enc", 6.15, 4.44, 0.90, 0.72, "weather · geo\n· time encoders", PURPLE, WHITE, 11.5)
    add_module(slide, "ML_C_fuse", 6.15, 5.47, 0.90, 0.72, "fusion", PURPLE, PURPLE_LIGHT, 13)
    add_module(slide, "ML_C_delta", 7.42, 4.88, 1.05, 1.10, "Δψ\nresidual MLP", PURPLE, PURPLE_LIGHT, 13)
    add_line(slide, "ML_C_u_enc", 5.78, 4.52, 6.15, 4.70, ORANGE, 1.0)
    add_line(slide, "ML_C_g_enc", 5.78, 5.66, 6.15, 4.90, GREEN, 1.0)
    add_line(slide, "ML_C_h_enc", 5.50, 6.68, 6.15, 5.06, MUTED, 1.0)
    add_line(slide, "ML_C_enc_fuse", 6.60, 5.16, 6.60, 5.47, PURPLE, 1.15)
    add_line(slide, "ML_C_fuse_delta", 7.05, 5.83, 7.42, 5.43, PURPLE, 1.15)
    add_text(slide, "ML_C_residual", 6.22, 6.45, 2.08, 0.42, "zₜ₊ₕ = zₜ + Δψ(zₜ,u,g,h)", 12.5, PURPLE, True)

    # (d) Zoom: explicit closure and multi-horizon output.
    add_compact_panel(slide, "ML_D", 8.96, yb, 4.21, hb, "(d) Explicit forecast closure", ORANGE_LIGHT, ORANGE)
    add_module(slide, "ML_D_O", 9.24, 4.26, 0.62, 0.90, "O", GREEN, GREEN_LIGHT, 18)
    add_module(slide, "ML_D_r", 10.15, 4.26, 0.70, 0.90, "rₕ", TEAL, TEAL_LIGHT, 16)
    add_plus(slide, "ML_D_plus", 11.13, 4.47, 0.46)
    add_line(slide, "ML_D_O_r", 9.86, 4.71, 10.15, 4.71, GREEN, 1.3)
    add_line(slide, "ML_D_r_plus", 10.85, 4.71, 11.13, 4.71, ORANGE, 1.3)
    add_module(slide, "ML_D_b", 9.52, 5.51, 1.58, 0.58, "history-only bₕ", BLUE, BLUE_LIGHT, 12)
    add_elbow(slide, "ML_D_b_plus", [(11.10, 5.80), (11.36, 5.80), (11.36, 4.93)], BLUE, 1.1)
    add_text(slide, "ML_D_eq", 11.78, 4.30, 1.10, 0.82, "ŷ(t+h)\n= bₕ+rₕ", 13, ORANGE, True)
    for idx, horizon in enumerate(("h=5", "h=10", "h=20")):
        add_image_slot(slide, f"ML_D_out_{idx}", 9.20 + idx * 1.20, 6.22, 1.00, 0.82, horizon, "NDVI", ORANGE)


def transplant_slide1(original: Path, generated: Path, destination: Path) -> None:
    """Replace only ppt/slides/slide1.xml; preserve every other package part."""

    with zipfile.ZipFile(generated, "r") as source_zip:
        slide_xml = source_zip.read("ppt/slides/slide1.xml")

    with zipfile.ZipFile(original, "r") as input_zip, zipfile.ZipFile(
        destination, "w", compression=zipfile.ZIP_DEFLATED
    ) as output_zip:
        for info in input_zip.infolist():
            data = input_zip.read(info.filename)
            if info.filename == "ppt/slides/slide1.xml":
                data = slide_xml
            output_zip.writestr(info, data)


def build_one(original: Path, destination: Path, builder) -> None:
    prs = Presentation(original)
    if not prs.slides:
        raise RuntimeError(f"No slides in {original}")
    builder(prs.slides[0])
    with tempfile.TemporaryDirectory(prefix="terrastate_ppt_") as temp_dir:
        generated = Path(temp_dir) / original.name
        prs.save(generated)
        transplant_slide1(original, generated, destination)


def build_options_deck(destination: Path, builders) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for builder in builders:
        slide = prs.slides.add_slide(blank)
        builder(slide)
    prs.save(destination)


def main() -> None:
    fig1 = EXAMPLE / "fig1.pptx"
    fig2 = EXAMPLE / "fig2.pptx"
    if not fig1.exists() or not fig2.exists():
        raise FileNotFoundError("Expected 示例/fig1.pptx and 示例/fig2.pptx")

    backup1 = EXAMPLE / "fig1_before_codex_20260727.pptx"
    backup2 = EXAMPLE / "fig2_before_codex_20260727.pptx"
    if not backup1.exists():
        shutil.copy2(fig1, backup1)
    if not backup2.exists():
        shutil.copy2(fig2, backup2)

    # Paper numbering is narrative-first:
    # Figure 1 introduces the conceptual gap and evidence hierarchy;
    # Figure 2 gives the concrete architecture and dataflow.
    draft1 = OUTDIR / "fig1_concept_v4.pptx"
    draft2 = OUTDIR / "fig2_method_v4.pptx"
    build_one(fig1, draft1, build_concept_slide_v2)
    build_one(fig2, draft2, build_method_slide_v2)
    options1 = OUTDIR / "fig1_concept_design_options.pptx"
    options2 = OUTDIR / "fig2_method_design_options.pptx"
    selected1 = OUTDIR / "fig1_selected_v5.pptx"
    selected2 = OUTDIR / "fig2_selected_v6.pptx"
    build_options_deck(options1, [build_concept_slide_v2, build_concept_slide_thesis])
    build_options_deck(
        options2,
        [build_method_slide_v2, build_method_slide_dense, build_method_slide_multilevel],
    )
    build_one(fig1, selected1, build_concept_slide_thesis)
    build_one(fig2, selected2, build_method_slide_multilevel)

    print(draft1)
    print(draft2)
    print(options1)
    print(options2)
    print(selected1)
    print(selected2)
    print(backup1)
    print(backup2)


if __name__ == "__main__":
    main()
