#!/usr/bin/env python3
"""Build the author-approved Figure 2 adaptation from editable vector primitives.

The friend-provided slide is used only as a visual-organization reference.
No raster asset from that deck is reused: several embedded images have unknown
provenance and could be mistaken for TerraState inputs or outputs.

Outputs:
  source/fig2_friend_adapted.pptx  native editable PowerPoint
  source/fig2_friend_adapted.svg   editable vector master
  export/fig2_friend_adapted.pdf   paper vector
  export/fig2_friend_adapted.png   300 dpi preview
  qa/fig2_friend_adapted_grayscale.png
  qa/fig2_friend_adapted_paperscale.pdf/png
"""

from __future__ import annotations

import hashlib
import json
import sys
from pathlib import Path
from xml.sax.saxutils import escape

import fitz


ROOT = Path(__file__).resolve().parents[1]
PROJECT = ROOT.parent
sys.path.insert(0, str(PROJECT / ".tools" / "python-pptx"))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.dml import MSO_LINE_DASH_STYLE  # noqa: E402
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE  # noqa: E402
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN  # noqa: E402
from pptx.oxml.xmlchemy import OxmlElement  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402


SOURCE = ROOT / "source"
EXPORT = ROOT / "export"
QA = ROOT / "qa"
PAPER_FIGURES = PROJECT / "paper" / "figures"

W_IN = 7.0
H_IN = 3.18
SCALE = 300.0

INK = "202833"
MUTED = "66717C"
LINE = "AEB7C0"
WHITE = "FFFFFF"
PANEL = "FCFDFE"
BLUE = "2F6F9F"
BLUE_LIGHT = "EAF3F8"
TEAL = "238B7B"
TEAL_LIGHT = "EAF6F3"
PURPLE = "7256A8"
PURPLE_LIGHT = "F2EFF9"
ORANGE = "C96A16"
ORANGE_LIGHT = "FFF3E8"
GREEN = "467A55"
GREEN_LIGHT = "EDF5EE"
GRAY_LIGHT = "F3F5F7"
FONT = "Arial"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def add_arrowhead(connector) -> None:
    line = connector._element.spPr.ln
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")
    line.append(tail)


def ppt_text(
    slide,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str,
    size: float = 9.0,
    color: str = INK,
    bold: bool = False,
    align=PP_ALIGN.CENTER,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.015)
    frame.margin_top = frame.margin_bottom = Inches(0.008)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return shape


def ppt_box(
    slide,
    name: str,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    stroke: str,
    text: str = "",
    size: float = 9.0,
    bold: bool = False,
    rounded: bool = True,
    dashed: bool = False,
):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(stroke)
    shape.line.width = Pt(0.85)
    if dashed:
        shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if text:
        frame = shape.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = frame.margin_right = Inches(0.015)
        frame.margin_top = frame.margin_bottom = Inches(0.008)
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        paragraph = frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        run = paragraph.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(INK)
    return shape


def ppt_line(
    slide,
    name: str,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = INK,
    width: float = 1.0,
    arrow: bool = True,
    dashed: bool = False,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.name = name
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    if dashed:
        line.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if arrow:
        add_arrowhead(line)
    return line


def ppt_panel(slide, prefix, x, w, title, tint, accent):
    ppt_box(slide, f"{prefix}_panel", x, 0.04, w, 3.02, PANEL, "D3D8DE")
    ppt_box(
        slide,
        f"{prefix}_header",
        x,
        0.04,
        w,
        0.30,
        tint,
        accent,
        title,
        8.6,
        True,
    )


def ppt_grid(slide, prefix, x, y, w, h, accent, changed=False):
    ppt_box(slide, prefix, x, y, w, h, TEAL_LIGHT, accent)
    cols = rows = 4
    gx, gy = x + 0.07, y + 0.07
    gw, gh = w - 0.14, h - 0.30
    cell_w, cell_h = gw / cols, gh / rows
    fills = (TEAL, "63A99E", "A7D0C9", "DCEDEA")
    for col in range(cols):
        for row in range(rows):
            fill = fills[(col + row + (1 if changed else 0)) % len(fills)]
            ppt_box(
                slide,
                f"{prefix}_{col}_{row}",
                gx + col * cell_w,
                gy + row * cell_h,
                cell_w - 0.008,
                cell_h - 0.008,
                fill,
                WHITE,
                rounded=False,
            )


def ppt_history_icon(slide, x, y):
    for index, (dx, dy, fill) in enumerate(
        ((0.00, 0.08, "D7E8D0"), (0.08, 0.04, "BDD7C1"), (0.16, 0.00, "9FC6B0"))
    ):
        ppt_box(
            slide,
            f"A_history_frame_{index}",
            x + dx,
            y + dy,
            0.55,
            0.34,
            fill,
            BLUE,
            rounded=False,
        )
        ppt_line(
            slide,
            f"A_history_field_{index}",
            x + dx + 0.05,
            y + dy + 0.17,
            x + dx + 0.50,
            y + dy + 0.17,
            WHITE,
            0.45,
            False,
        )
    ppt_line(slide, "A_cloud_mask", x + 0.18, y + 0.03, x + 0.64, y + 0.36, BLUE, 0.65, False, True)


def ppt_weather_strip(slide, prefix, x, y, w, h, label):
    ppt_box(slide, prefix, x, y, w, h, PURPLE_LIGHT, PURPLE)
    colors = (ORANGE, BLUE, PURPLE)
    for row, color in enumerate(colors):
        yy = y + 0.08 + row * 0.12
        points = [
            (x + 0.08, yy + 0.04),
            (x + w * 0.34, yy + (0.00 if row == 1 else 0.08)),
            (x + w * 0.63, yy + (0.08 if row == 0 else 0.02)),
            (x + w - 0.08, yy + 0.05),
        ]
        for idx, (start, end) in enumerate(zip(points, points[1:])):
            ppt_line(
                slide,
                f"{prefix}_line_{row}_{idx}",
                *start,
                *end,
                color,
                0.55,
                False,
            )
    ppt_text(slide, f"{prefix}_label", x + 0.04, y + h - 0.22, w - 0.08, 0.18, label, 7.8, INK, True)


def build_pptx(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(WHITE)

    ppt_panel(slide, "A", 0.04, 1.24, "(a) Historical context", BLUE_LIGHT, BLUE)
    ppt_panel(slide, "B", 1.32, 1.84, "(b) Predictive-state construction", TEAL_LIGHT, TEAL)
    ppt_panel(slide, "C", 3.20, 2.12, "(c) Weather-conditioned dynamics", PURPLE_LIGHT, PURPLE)
    ppt_panel(slide, "D", 5.36, 1.60, "(d) Forecast closure", ORANGE_LIGHT, ORANGE)

    # Panel A: history-only inputs. All visual elements are schematic.
    ppt_box(slide, "A_eo_card", 0.14, 0.47, 1.04, 0.66, WHITE, BLUE)
    ppt_history_icon(slide, 0.28, 0.58)
    ppt_text(slide, "A_eo_label", 0.18, 0.88, 0.96, 0.22, "cloud-masked\nEO history", 9.0, INK, True)
    ppt_box(slide, "A_met_card", 0.14, 1.30, 1.04, 0.55, BLUE_LIGHT, BLUE)
    for idx, yy in enumerate((1.43, 1.57)):
        ppt_line(slide, f"A_met_{idx}_1", 0.25, yy, 0.51, yy - 0.05, BLUE, 0.55, False)
        ppt_line(slide, f"A_met_{idx}_2", 0.51, yy - 0.05, 0.78, yy + 0.04, BLUE, 0.55, False)
        ppt_line(slide, f"A_met_{idx}_3", 0.78, yy + 0.04, 1.06, yy - 0.02, BLUE, 0.55, False)
    ppt_text(slide, "A_met_label", 0.18, 1.64, 0.96, 0.18, "past meteorology", 9.0, INK, True)
    ppt_box(slide, "A_geo_card", 0.14, 2.02, 1.04, 0.62, GREEN_LIGHT, GREEN)
    ppt_box(slide, "A_landcover", 0.25, 2.13, 0.35, 0.27, "C9DDB9", GREEN, rounded=False)
    ppt_box(slide, "A_dem", 0.70, 2.13, 0.35, 0.27, "E4E7EA", MUTED, rounded=False)
    for idx in range(3):
        ppt_line(slide, f"A_contour_{idx}", 0.73, 2.18 + idx * 0.06, 1.02, 2.15 + idx * 0.08, MUTED, 0.35, False)
    ppt_text(slide, "A_geo_label", 0.18, 2.43, 0.96, 0.16, "static geography  g", 9.0, INK, True)

    # Panel B: one history encoder, state branch, and context-only branch.
    ppt_box(slide, "B_q", 1.43, 0.69, 0.62, 0.70, BLUE_LIGHT, BLUE, "history\nencoder qθ", 8.7, True)
    for col in range(3):
        for row in range(3):
            ppt_box(
                slide,
                f"B_token_{col}_{row}",
                2.16 + col * 0.10,
                0.82 + row * 0.10,
                0.075,
                0.075,
                (BLUE, TEAL, PURPLE)[col],
                WHITE,
                rounded=False,
            )
    ppt_box(slide, "B_P", 2.61, 0.76, 0.32, 0.54, TEAL_LIGHT, TEAL, "Pρ", 10.0, True)
    ppt_grid(slide, "B_zt", 2.97, 0.66, 0.43, 0.78, TEAL)
    ppt_text(slide, "B_zt_label", 2.97, 1.18, 0.43, 0.17, "zₜ", 9.0, TEAL, True)
    ppt_line(slide, "A_to_q", 1.18, 1.04, 1.43, 1.04, BLUE, 1.15)
    ppt_line(slide, "B_q_to_tokens", 2.05, 1.04, 2.16, 1.04, BLUE, 1.0)
    ppt_line(slide, "B_tokens_to_P", 2.49, 1.04, 2.61, 1.04, TEAL, 1.0)
    ppt_line(slide, "B_P_to_zt", 2.93, 1.04, 2.97, 1.04, TEAL, 1.0)
    ppt_box(slide, "B_bh", 1.68, 2.12, 1.15, 0.52, BLUE_LIGHT, BLUE, "context-only forecast\nbₕ", 9.0, True)
    ppt_line(slide, "B_q_to_bh_1", 1.74, 1.39, 1.74, 2.18, BLUE, 0.9, False)
    ppt_line(slide, "B_q_to_bh_2", 1.74, 2.18, 1.90, 2.18, BLUE, 0.9)

    # Panel C: future weather is introduced only at the shared transition.
    ppt_weather_strip(slide, "C_weather", 3.31, 0.47, 0.94, 0.60, "future weather  u")
    ppt_box(slide, "C_g", 4.32, 0.53, 0.42, 0.48, GREEN_LIGHT, GREEN, "static\ng", 8.2, True)
    ppt_box(slide, "C_h", 4.82, 0.53, 0.38, 0.48, GRAY_LIGHT, MUTED, "horizon\nh", 8.0, True)
    ppt_box(slide, "C_T", 3.79, 1.31, 0.92, 0.77, PURPLE_LIGHT, PURPLE, "shared transition\nTψ", 9.2, True)
    ppt_grid(slide, "C_zfuture", 4.84, 1.31, 0.40, 0.77, TEAL, True)
    ppt_text(slide, "C_zfuture_label", 4.84, 1.84, 0.40, 0.18, "zₜ₊ₕ", 9.0, TEAL, True)
    ppt_line(slide, "B_zt_to_T", 3.40, 1.04, 3.79, 1.68, TEAL, 1.15)
    ppt_line(slide, "C_weather_to_T", 3.78, 1.07, 4.03, 1.31, PURPLE, 0.95)
    ppt_line(slide, "C_g_to_T", 4.53, 1.01, 4.46, 1.31, GREEN, 0.85)
    ppt_line(slide, "C_h_to_T", 5.01, 1.01, 4.66, 1.31, MUTED, 0.85)
    ppt_line(slide, "C_T_to_zfuture", 4.71, 1.69, 4.84, 1.69, TEAL, 1.15)
    ppt_box(slide, "C_Q3", 3.34, 2.35, 1.08, 0.34, WHITE, PURPLE, "Q3  replace u", 9.0, True, dashed=True)
    ppt_line(slide, "C_Q3_pointer", 3.88, 2.35, 3.77, 1.08, PURPLE, 0.75, True, True)
    ppt_box(slide, "C_T_identity", 4.57, 2.35, 0.63, 0.34, WHITE, ORANGE, "T→I\nsupport", 8.5, True, dashed=True)
    ppt_line(slide, "C_T_identity_pointer", 4.88, 2.35, 4.56, 2.08, ORANGE, 0.70, True, True)

    # Panel D: state contribution and context forecast close explicitly.
    ppt_text(slide, "D_readout_label", 5.45, 0.48, 0.55, 0.18, "state readout", 9.0, GREEN, True)
    ppt_box(slide, "D_O", 5.52, 0.75, 0.36, 0.55, GREEN_LIGHT, GREEN, "Oω", 10.0, True)
    ppt_box(slide, "D_rh", 5.98, 0.75, 0.36, 0.55, TEAL_LIGHT, TEAL, "rₕ", 10.2, True)
    plus = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.48), Inches(0.83), Inches(0.31), Inches(0.31))
    plus.name = "D_plus"
    plus.fill.solid()
    plus.fill.fore_color.rgb = rgb(WHITE)
    plus.line.color.rgb = rgb(ORANGE)
    plus.line.width = Pt(1.0)
    frame = plus.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.add_run()
    run.text = "+"
    run.font.name = FONT
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = rgb(ORANGE)
    ppt_line(slide, "C_zfuture_to_O", 5.24, 1.69, 5.52, 1.03, TEAL, 1.10)
    ppt_line(slide, "D_O_to_rh", 5.88, 1.03, 5.98, 1.03, GREEN, 1.0)
    ppt_line(slide, "D_rh_to_plus", 6.34, 1.03, 6.48, 0.99, ORANGE, 1.0)
    ppt_box(slide, "D_forecast", 5.63, 1.88, 1.12, 0.62, ORANGE_LIGHT, ORANGE)
    for col in range(4):
        for row in range(3):
            fill = ("D8E8C9", "B9D4B5", "91BC9B", "E6D7A9")[(col + row) % 4]
            ppt_box(
                slide,
                f"D_forecast_cell_{col}_{row}",
                5.78 + col * 0.16,
                2.00 + row * 0.11,
                0.145,
                0.105,
                fill,
                WHITE,
                rounded=False,
            )
    ppt_text(slide, "D_forecast_label", 5.68, 2.31, 1.02, 0.18, "forecast  ŷₜ₊ₕ", 9.0, ORANGE, True)
    ppt_line(slide, "D_plus_to_forecast", 6.64, 1.14, 6.32, 1.88, ORANGE, 1.0)
    ppt_line(slide, "D_bh_path_1", 2.83, 2.38, 3.03, 2.38, BLUE, 0.95, False)
    ppt_line(slide, "D_bh_path_2", 3.03, 2.38, 3.03, 2.84, BLUE, 0.95, False)
    ppt_line(slide, "D_bh_path_3", 3.03, 2.84, 6.84, 2.84, BLUE, 0.95, False)
    ppt_line(slide, "D_bh_path_4", 6.84, 2.84, 6.84, 1.14, BLUE, 0.95)
    ppt_box(slide, "D_Q2", 5.62, 1.42, 1.08, 0.34, WHITE, ORANGE, "Q2  remove rₕ", 9.0, True, dashed=True)
    ppt_line(slide, "D_Q2_pointer", 6.16, 1.42, 6.16, 1.30, ORANGE, 0.75, True, True)

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(path)


def svg_rect(x, y, w, h, fill, stroke, rx=12, dash=""):
    dash_attr = f' stroke-dasharray="{dash}"' if dash else ""
    return (
        f'<rect x="{x*SCALE:.1f}" y="{y*SCALE:.1f}" width="{w*SCALE:.1f}" '
        f'height="{h*SCALE:.1f}" rx="{rx}" fill="#{fill}" stroke="#{stroke}" '
        f'stroke-width="4"{dash_attr}/>'
    )


def svg_text(x, y, w, h, value, size=9.0, color=INK, bold=False, anchor="middle"):
    font_px = size * 4.1667
    weight = "700" if bold else "400"
    lines = value.split("\n")
    line_h = font_px * 1.08
    center_y = (y + h / 2) * SCALE
    start_y = center_y - (len(lines) - 1) * line_h / 2
    if anchor == "middle":
        tx = (x + w / 2) * SCALE
    else:
        tx = x * SCALE
    spans = "".join(
        f'<tspan x="{tx:.1f}" y="{start_y + idx*line_h:.1f}">{escape(line)}</tspan>'
        for idx, line in enumerate(lines)
    )
    return (
        f'<text text-anchor="{anchor}" font-family="Arial,Helvetica,sans-serif" '
        f'font-size="{font_px:.1f}" font-weight="{weight}" fill="#{color}">{spans}</text>'
    )


def svg_line(x1, y1, x2, y2, color=INK, width=1.0, arrow=True, dashed=False):
    marker = ' marker-end="url(#arrow)"' if arrow else ""
    dash = ' stroke-dasharray="12 9"' if dashed else ""
    return (
        f'<line x1="{x1*SCALE:.1f}" y1="{y1*SCALE:.1f}" '
        f'x2="{x2*SCALE:.1f}" y2="{y2*SCALE:.1f}" '
        f'stroke="#{color}" stroke-width="{width*4:.1f}"{dash}{marker}/>'
    )


def svg_panel(parts, x, w, title, tint, accent):
    parts.append(svg_rect(x, 0.04, w, 3.02, PANEL, "D3D8DE", 18))
    parts.append(svg_rect(x, 0.04, w, 0.30, tint, accent, 18))
    parts.append(svg_text(x + 0.02, 0.04, w - 0.04, 0.30, title, 8.6, INK, True))


def svg_grid(parts, prefix, x, y, w, h, changed=False):
    del prefix
    parts.append(svg_rect(x, y, w, h, TEAL_LIGHT, TEAL, 12))
    fills = (TEAL, "63A99E", "A7D0C9", "DCEDEA")
    gx, gy = x + 0.07, y + 0.07
    cell_w, cell_h = (w - 0.14) / 4, (h - 0.30) / 4
    for col in range(4):
        for row in range(4):
            fill = fills[(col + row + (1 if changed else 0)) % len(fills)]
            parts.append(svg_rect(gx + col*cell_w, gy + row*cell_h, cell_w - 0.008, cell_h - 0.008, fill, WHITE, 0))


def build_svg(path: Path) -> None:
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{W_IN*SCALE:.0f}" '
        f'height="{H_IN*SCALE:.0f}" viewBox="0 0 {W_IN*SCALE:.0f} {H_IN*SCALE:.0f}">',
        "<defs><marker id=\"arrow\" markerWidth=\"9\" markerHeight=\"9\" refX=\"8\" refY=\"3\" "
        "orient=\"auto\" markerUnits=\"strokeWidth\"><path d=\"M0,0 L0,6 L9,3 z\" "
        f"fill=\"#{INK}\"/></marker></defs>",
        f'<rect width="100%" height="100%" fill="#{WHITE}"/>',
    ]
    svg_panel(parts, 0.04, 1.24, "(a) Historical context", BLUE_LIGHT, BLUE)
    svg_panel(parts, 1.32, 1.84, "(b) Predictive-state construction", TEAL_LIGHT, TEAL)
    svg_panel(parts, 3.20, 2.12, "(c) Weather-conditioned dynamics", PURPLE_LIGHT, PURPLE)
    svg_panel(parts, 5.36, 1.60, "(d) Forecast closure", ORANGE_LIGHT, ORANGE)

    # A
    parts.append(svg_rect(0.14, 0.47, 1.04, 0.66, WHITE, BLUE))
    for dx, dy, fill in ((0, .08, "D7E8D0"), (.08, .04, "BDD7C1"), (.16, 0, "9FC6B0")):
        parts.append(svg_rect(.28+dx, .58+dy, .55, .34, fill, BLUE, 0))
        parts.append(svg_line(.33+dx, .75+dy, .78+dx, .75+dy, WHITE, .45, False))
    parts.append(svg_line(.46, .61, .92, .94, BLUE, .65, False, True))
    parts.append(svg_text(.18, .88, .96, .22, "cloud-masked\nEO history", 9.0, INK, True))
    parts.append(svg_rect(.14, 1.30, 1.04, .55, BLUE_LIGHT, BLUE))
    for row, yy in enumerate((1.43, 1.57)):
        del row
        parts.extend([
            svg_line(.25, yy, .51, yy-.05, BLUE, .55, False),
            svg_line(.51, yy-.05, .78, yy+.04, BLUE, .55, False),
            svg_line(.78, yy+.04, 1.06, yy-.02, BLUE, .55, False),
        ])
    parts.append(svg_text(.18, 1.64, .96, .18, "past meteorology", 9.0, INK, True))
    parts.append(svg_rect(.14, 2.02, 1.04, .62, GREEN_LIGHT, GREEN))
    parts.append(svg_rect(.25, 2.13, .35, .27, "C9DDB9", GREEN, 0))
    parts.append(svg_rect(.70, 2.13, .35, .27, "E4E7EA", MUTED, 0))
    for idx in range(3):
        parts.append(svg_line(.73, 2.18+idx*.06, 1.02, 2.15+idx*.08, MUTED, .35, False))
    parts.append(svg_text(.18, 2.43, .96, .16, "static geography  g", 9.0, INK, True))

    # B
    parts.extend([
        svg_rect(1.43, .69, .62, .70, BLUE_LIGHT, BLUE),
        svg_text(1.43, .69, .62, .70, "history\nencoder qθ", 8.7, INK, True),
    ])
    for col in range(3):
        for row in range(3):
            parts.append(svg_rect(2.16+col*.10, .82+row*.10, .075, .075, (BLUE, TEAL, PURPLE)[col], WHITE, 0))
    parts.extend([
        svg_rect(2.61, .76, .32, .54, TEAL_LIGHT, TEAL),
        svg_text(2.61, .76, .32, .54, "Pρ", 10, INK, True),
    ])
    svg_grid(parts, "B_zt", 2.97, .66, .43, .78)
    parts.append(svg_text(2.97, 1.18, .43, .17, "zₜ", 9.0, TEAL, True))
    parts.extend([
        svg_line(1.18, 1.04, 1.43, 1.04, BLUE, 1.15),
        svg_line(2.05, 1.04, 2.16, 1.04, BLUE, 1.0),
        svg_line(2.49, 1.04, 2.61, 1.04, TEAL, 1.0),
        svg_line(2.93, 1.04, 2.97, 1.04, TEAL, 1.0),
        svg_rect(1.68, 2.12, 1.15, .52, BLUE_LIGHT, BLUE),
        svg_text(1.68, 2.12, 1.15, .52, "context-only forecast\nbₕ", 9.0, INK, True),
        svg_line(1.74, 1.39, 1.74, 2.18, BLUE, .9, False),
        svg_line(1.74, 2.18, 1.90, 2.18, BLUE, .9),
    ])

    # C
    parts.append(svg_rect(3.31, .47, .94, .60, PURPLE_LIGHT, PURPLE))
    for row, color in enumerate((ORANGE, BLUE, PURPLE)):
        yy = .55 + row*.12
        p = [(3.39, yy+.04), (3.63, yy), (3.88, yy+.07), (4.17, yy+.03)]
        for a, b in zip(p, p[1:]):
            parts.append(svg_line(*a, *b, color, .55, False))
    parts.append(svg_text(3.35, .84, .86, .18, "future weather  u", 9.0, INK, True))
    parts.extend([
        svg_rect(4.32, .53, .42, .48, GREEN_LIGHT, GREEN),
        svg_text(4.32, .53, .42, .48, "static\ng", 8.2, INK, True),
        svg_rect(4.82, .53, .38, .48, GRAY_LIGHT, MUTED),
        svg_text(4.82, .53, .38, .48, "horizon\nh", 8.0, INK, True),
        svg_rect(3.79, 1.31, .92, .77, PURPLE_LIGHT, PURPLE),
        svg_text(3.79, 1.31, .92, .77, "shared transition\nTψ", 9.2, INK, True),
    ])
    svg_grid(parts, "C_zfuture", 4.84, 1.31, .40, .77, True)
    parts.extend([
        svg_text(4.84, 1.84, .40, .18, "zₜ₊ₕ", 9.0, TEAL, True),
        svg_line(3.40, 1.04, 3.79, 1.68, TEAL, 1.15),
        svg_line(3.78, 1.07, 4.03, 1.31, PURPLE, .95),
        svg_line(4.53, 1.01, 4.46, 1.31, GREEN, .85),
        svg_line(5.01, 1.01, 4.66, 1.31, MUTED, .85),
        svg_line(4.71, 1.69, 4.84, 1.69, TEAL, 1.15),
        svg_rect(3.34, 2.35, 1.08, .34, WHITE, PURPLE, 12, "12 9"),
        svg_text(3.34, 2.35, 1.08, .34, "Q3  replace u", 9.0, INK, True),
        svg_line(3.88, 2.35, 3.77, 1.08, PURPLE, .75, True, True),
        svg_rect(4.57, 2.35, .63, .34, WHITE, ORANGE, 12, "12 9"),
        svg_text(4.57, 2.35, .63, .34, "T→I\nsupport", 8.5, INK, True),
        svg_line(4.88, 2.35, 4.56, 2.08, ORANGE, .70, True, True),
    ])

    # D
    parts.extend([
        svg_text(5.45, .48, .55, .18, "state readout", 9.0, GREEN, True),
        svg_rect(5.52, .75, .36, .55, GREEN_LIGHT, GREEN),
        svg_text(5.52, .75, .36, .55, "Oω", 10.0, INK, True),
        svg_rect(5.98, .75, .36, .55, TEAL_LIGHT, TEAL),
        svg_text(5.98, .75, .36, .55, "rₕ", 10.2, INK, True),
        f'<circle cx="{6.635*SCALE:.1f}" cy="{.985*SCALE:.1f}" r="{.155*SCALE:.1f}" fill="#{WHITE}" stroke="#{ORANGE}" stroke-width="4"/>',
        svg_text(6.48, .83, .31, .31, "+", 11.0, ORANGE, True),
        svg_line(5.24, 1.69, 5.52, 1.03, TEAL, 1.10),
        svg_line(5.88, 1.03, 5.98, 1.03, GREEN, 1.0),
        svg_line(6.34, 1.03, 6.48, .99, ORANGE, 1.0),
        svg_rect(5.63, 1.88, 1.12, .62, ORANGE_LIGHT, ORANGE),
    ])
    fills = ("D8E8C9", "B9D4B5", "91BC9B", "E6D7A9")
    for col in range(4):
        for row in range(3):
            parts.append(svg_rect(5.78+col*.16, 2.00+row*.11, .145, .095, fills[(col+row)%4], WHITE, 0))
    parts.extend([
        svg_text(5.68, 2.31, 1.02, .18, "forecast  ŷₜ₊ₕ", 9.0, ORANGE, True),
        svg_line(6.64, 1.14, 6.32, 1.88, ORANGE, 1.0),
        svg_line(2.83, 2.38, 3.03, 2.38, BLUE, .95, False),
        svg_line(3.03, 2.38, 3.03, 2.84, BLUE, .95, False),
        svg_line(3.03, 2.84, 6.84, 2.84, BLUE, .95, False),
        svg_line(6.84, 2.84, 6.84, 1.14, BLUE, .95),
        svg_rect(5.62, 1.42, 1.08, .34, WHITE, ORANGE, 12, "12 9"),
        svg_text(5.62, 1.42, 1.08, .34, "Q2  remove rₕ", 9.0, INK, True),
        svg_line(6.16, 1.42, 6.16, 1.30, ORANGE, .75, True, True),
        "</svg>",
    ])
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text("\n".join(parts), encoding="utf-8")


def export(svg: Path, pdf: Path, png: Path, gray: Path, preview_pdf: Path, preview_png: Path) -> None:
    document = fitz.open(svg)
    pdf_bytes = document.convert_to_pdf()
    vector = fitz.open("pdf", pdf_bytes)
    vector.save(pdf, garbage=4, deflate=True)
    page = vector[0]
    pixmap = page.get_pixmap(dpi=300, alpha=False)
    pixmap.save(png)
    fitz.Pixmap(fitz.csGRAY, pixmap).save(gray)

    preview = fitz.open()
    sheet = preview.new_page(width=612, height=792)
    sheet.insert_text((54, 46), "Figure 2 at AAAI full width (7.0 in)", fontsize=9)
    rect = fitz.Rect(54, 70, 558, 70 + H_IN * 72)
    sheet.show_pdf_page(rect, vector, 0)
    preview.save(preview_pdf, garbage=4, deflate=True)
    sheet.get_pixmap(dpi=180, alpha=False).save(preview_png)


def main() -> None:
    SOURCE.mkdir(parents=True, exist_ok=True)
    EXPORT.mkdir(parents=True, exist_ok=True)
    QA.mkdir(parents=True, exist_ok=True)
    PAPER_FIGURES.mkdir(parents=True, exist_ok=True)

    pptx = SOURCE / "fig2_friend_adapted.pptx"
    svg = SOURCE / "fig2_friend_adapted.svg"
    pdf = EXPORT / "fig2_friend_adapted.pdf"
    png = EXPORT / "fig2_friend_adapted.png"
    gray = QA / "fig2_friend_adapted_grayscale.png"
    preview_pdf = QA / "fig2_friend_adapted_paperscale.pdf"
    preview_png = QA / "fig2_friend_adapted_paperscale.png"

    build_pptx(pptx)
    build_svg(svg)
    export(svg, pdf, png, gray, preview_pdf, preview_png)

    paper_pdf = PAPER_FIGURES / "terrastate_architecture_fig2.pdf"
    paper_svg = PAPER_FIGURES / "terrastate_architecture_fig2.svg"
    paper_png = PAPER_FIGURES / "terrastate_architecture_fig2.png"
    paper_pptx = PAPER_FIGURES / "terrastate_architecture_fig2.pptx"
    for source, destination in (
        (pdf, paper_pdf),
        (svg, paper_svg),
        (png, paper_png),
        (pptx, paper_pptx),
    ):
        destination.write_bytes(source.read_bytes())

    manifest = {
        "source_reference": "示例/fig2——2.pptx slide 1",
        "reuse_policy": "layout principles only; no embedded raster asset reused",
        "editable_pptx": str(pptx.relative_to(ROOT)),
        "editable_svg": str(svg.relative_to(ROOT)),
        "vector_pdf": str(pdf.relative_to(ROOT)),
        "png_300dpi": str(png.relative_to(ROOT)),
        "grayscale": str(gray.relative_to(ROOT)),
        "paperscale_pdf": str(preview_pdf.relative_to(ROOT)),
        "paperscale_png": str(preview_png.relative_to(ROOT)),
        "paper_asset": str(paper_pdf.relative_to(PROJECT)),
        "dimensions_inches": [W_IN, H_IN],
        "sha256": {
            path.name: sha256(path)
            for path in (pptx, svg, pdf, png, gray, preview_pdf, preview_png)
        },
    }
    (ROOT / "FIG2_FRIEND_ADAPTATION_MANIFEST.json").write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False) + "\n",
        encoding="utf-8",
    )
    print(json.dumps(manifest, indent=2, ensure_ascii=False))


if __name__ == "__main__":
    main()
